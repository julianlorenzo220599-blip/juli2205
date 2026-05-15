"""CLI end-to-end: factura + ubicación + PDI → revisión interna + propuesta cliente.

Uso (ejemplo manual):
    py -m rv_propuestas.cli demo
    py -m rv_propuestas.cli desde-factura --pdf factura.pdf --lat -34.6 --lon -58.4 \
        --tension-pdi 0.38 --fases 3 --capacidad-pdi 300 --salida ./output

Uso (programático): ver rv_propuestas.pipeline.ejecutar_pipeline()
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path
from typing import Optional

from .inputs.facturas import Factura, from_manual, parse_pdf
from .inputs.pdi import desde_factura_y_pdi
from .inputs.ubicacion import Ubicacion, get_irradiacion, estimar_offline
from .pipeline import ejecutar_pipeline


def cmd_demo(args) -> None:
    """Corre un caso demo con datos sintéticos para validar el pipeline."""
    factura = from_manual(
        distribuidora="EDENOR",
        consumos_kwh=[42000, 39000, 38000, 35000, 31000, 28000,
                      29000, 31000, 34000, 37000, 39000, 41000],
        potencia_contratada_kw=150.0,
        tension_suministro="BT 380V",
        categoria_tarifaria="T3-MD",
        titular="Cliente Demo S.A.",
        nis="12345678",
    )
    ubic = Ubicacion(nombre="Buenos Aires", lat=-34.6037, lon=-58.3816)
    pdi = desde_factura_y_pdi(tension_kv=0.38, fases=3, capacidad_disponible_kw=200.0,
                              distancia_al_tablero_m=20.0)
    salida = Path(args.salida or "./output")
    ejecutar_pipeline(
        factura=factura,
        ubicacion=ubic,
        pdi=pdi,
        salida=salida,
        cliente_nombre="Cliente Demo S.A.",
        proyecto_nombre="Demo 200 kW",
        usar_pvgis=not args.offline,
        precios_path=args.precios,
        template_ppt=args.template,
        pvsyst_memo=args.pvsyst_memo,
        pvsyst_report=args.pvsyst_report,
        clickup_push=args.clickup,
        clickup_list_id=args.clickup_list,
    )


def cmd_desde_factura(args) -> None:
    factura = parse_pdf(args.pdf)
    ubic = Ubicacion(nombre=args.ubicacion or "Sin nombre", lat=args.lat, lon=args.lon)
    pdi = desde_factura_y_pdi(
        tension_kv=args.tension_pdi,
        fases=args.fases,
        capacidad_disponible_kw=args.capacidad_pdi,
        distancia_al_tablero_m=args.distancia or 10.0,
        distribuidora=factura.distribuidora,
    )
    ejecutar_pipeline(
        factura=factura,
        ubicacion=ubic,
        pdi=pdi,
        salida=Path(args.salida or "./output"),
        cliente_nombre=args.cliente or factura.titular,
        proyecto_nombre=args.proyecto or "",
        usar_pvgis=not args.offline,
        precios_path=args.precios,
        template_ppt=args.template,
        pvsyst_memo=args.pvsyst_memo,
        pvsyst_report=args.pvsyst_report,
        clickup_push=args.clickup,
        clickup_list_id=args.clickup_list,
    )


def cmd_placeholders(args) -> None:
    """Inspecciona un .pptx y lista los placeholders {{...}} que contiene.

    Útil para que el diseñador valide su template antes de mandarlo a la
    pipeline. Compara contra el set de claves soportadas.
    """
    from pptx import Presentation
    from .render.template import listar_placeholders, contexto_de_propuesta

    prs = Presentation(args.template)
    detectados = listar_placeholders(prs)
    if not detectados:
        print(f"⚠ {args.template} no tiene placeholders {{...}} — modo programático")
        return

    # Calculamos las claves soportadas creando un contexto vacío. Como la firma
    # requiere objetos, usamos un dict de referencia hardcoded acá. Si se agregan
    # claves nuevas en contexto_de_propuesta, se reflejan automáticamente.
    soportadas = _claves_soportadas()
    print(f"Placeholders detectados en {args.template}:")
    for k in sorted(detectados):
        marker = "✓" if k in soportadas else "✗ DESCONOCIDO"
        print(f"  {marker}  {{{{{k}}}}}")

    no_usados = soportadas - detectados
    if no_usados and args.verbose:
        print(f"\nDisponibles pero no usados ({len(no_usados)}):")
        for k in sorted(no_usados):
            print(f"     {{{{{k}}}}}")


def _claves_soportadas() -> set[str]:
    """Llama a contexto_de_propuesta() con objetos sentinela para extraer las claves."""
    from types import SimpleNamespace
    from .render.template import contexto_de_propuesta
    sentinel_factura = SimpleNamespace(
        titular="", distribuidora="", categoria_tarifaria="",
        tension_suministro="", potencia_contratada_kw=0,
        consumo_anual_kwh=0, consumo_mensual_promedio=0,
        direccion="", nis="",
    )
    sentinel_sizing = SimpleNamespace(
        kwp_real=0, n_paneles=0, generacion_anual_kwh=0, cobertura=0,
    )
    sentinel_inv = SimpleNamespace(
        cantidad=0,
        inversor=SimpleNamespace(sku="", descripcion=""),
    )
    sentinel_costeo = SimpleNamespace(
        neto_cliente=0, iva_total=0, total_cliente=0,
    )
    ctx = contexto_de_propuesta(
        factura=sentinel_factura, sizing=sentinel_sizing,
        inv_cfg=sentinel_inv, costeo=sentinel_costeo,
    )
    return set(ctx.keys())


def main(argv: Optional[list[str]] = None) -> int:
    parser = argparse.ArgumentParser(prog="rv_propuestas", description="Generador de propuestas RV Energía")
    sub = parser.add_subparsers(dest="cmd", required=True)

    p_demo = sub.add_parser("demo", help="Corre un caso demo sintético")
    p_demo.add_argument("--salida", default="./output")
    p_demo.add_argument("--offline", action="store_true", help="Usar irradiación offline (sin PVGIS)")
    p_demo.add_argument("--precios", default="./data/precios.example.yaml")
    p_demo.add_argument("--template", default=None,
                        help="Ruta a un .pptx de RV con placeholders {{clave}}")
    p_demo.add_argument("--pvsyst-memo", action="store_true",
                        help="Generar memo PVSYST_INPUT_*.txt junto al resto del output")
    p_demo.add_argument("--pvsyst-report", default=None,
                        help="Path al CSV de PVSyst para override de generación")
    p_demo.add_argument("--clickup", action="store_true",
                        help="Push a ClickUp al finalizar (CLICKUP_API_TOKEN + LIST_ID)")
    p_demo.add_argument("--clickup-list", default=None)
    p_demo.set_defaults(func=cmd_demo)

    p_real = sub.add_parser("desde-factura", help="Procesa una factura PDF real")
    p_real.add_argument("--pdf", required=True, help="Ruta al PDF de factura")
    p_real.add_argument("--lat", type=float, required=True)
    p_real.add_argument("--lon", type=float, required=True)
    p_real.add_argument("--ubicacion", default="")
    p_real.add_argument("--tension-pdi", type=float, required=True, help="kV (0.38 para BT, 13.2/33 para MT)")
    p_real.add_argument("--fases", type=int, default=3)
    p_real.add_argument("--capacidad-pdi", type=float, default=None)
    p_real.add_argument("--distancia", type=float, default=None)
    p_real.add_argument("--cliente", default="")
    p_real.add_argument("--proyecto", default="")
    p_real.add_argument("--salida", default="./output")
    p_real.add_argument("--offline", action="store_true")
    p_real.add_argument("--precios", default="./data/precios.example.yaml")
    p_real.add_argument("--template", default=None,
                        help="Ruta a un .pptx de RV con placeholders {{clave}}")
    p_real.add_argument("--pvsyst-memo", action="store_true",
                        help="Generar PVSYST_INPUT_<proy>.txt con todos los parámetros "
                             "para transcribir a PVSyst")
    p_real.add_argument("--pvsyst-report", default=None,
                        help="Path al CSV 'Main results, per month' de PVSyst — "
                             "reemplaza la estimación interna por la validada")
    p_real.add_argument("--clickup", action="store_true",
                        help="Crear task automática en ClickUp con resumen + "
                             "adjuntos (requiere CLICKUP_API_TOKEN + CLICKUP_LIST_ID)")
    p_real.add_argument("--clickup-list", default=None,
                        help="Override del CLICKUP_LIST_ID de entorno")
    p_real.set_defaults(func=cmd_desde_factura)

    p_ph = sub.add_parser(
        "placeholders",
        help="Lista los placeholders {{...}} de un template .pptx",
    )
    p_ph.add_argument("--template", required=True, help="Ruta al .pptx a inspeccionar")
    p_ph.add_argument("-v", "--verbose", action="store_true",
                      help="Mostrar también claves disponibles pero no usadas")
    p_ph.set_defaults(func=cmd_placeholders)

    args = parser.parse_args(argv)
    args.func(args)
    return 0


if __name__ == "__main__":
    sys.exit(main())
