"""Sustitución de placeholders {{clave}} en templates PPT de RV.

Flujo recomendado para el equipo de marketing:
  1. Diseñador arma un .pptx con la identidad visual de RV.
  2. En los textos donde van datos del proyecto, escribe `{{clave}}` o
     `{{clave|filtro}}` (ej: `{{kwp|1}} kWp`, `Inversión total: {{total_usd|usd}}`).
  3. Pasamos el .pptx como `--template` al CLI: la pipeline rellena los
     placeholders preservando layout, colores y tipografía del template.

Si el template no tiene placeholders, el renderer cae al modo programático
(comportamiento legacy: agregar 5 slides de propuesta sobre el template).

Para descubrir qué placeholders aceptamos hoy, ver `contexto_de_propuesta()`
o correr `py -m rv_propuestas.cli placeholders --template foo.pptx`.
"""
from __future__ import annotations

import datetime
import re
from typing import Any

from ..config import fmt_ar

_PLACEHOLDER_RE = re.compile(
    r"\{\{\s*([a-zA-Z_][a-zA-Z0-9_]*)\s*(?:\|\s*([^}]*?)\s*)?\}\}"
)

# Marcadores de chart: `{{chart:consumo_mensual}}` — los reemplaza un chart real
_CHART_RE = re.compile(r"\{\{\s*chart\s*:\s*([a-zA-Z_][a-zA-Z0-9_]*)\s*\}\}")


def _formatear(value: Any, fmt: str | None) -> str:
    """Aplica el filtro `|fmt` al valor.

    Filtros soportados:
      ─ sin filtro     → smart (entero si lo es, sino 2 decimales)
      ─ `0`, `1`, `2`  → N decimales en formato AR
      ─ `usd`          → "USD 1.234.567"
      ─ `pct`          → "95%" (asume value en [0..1])
      ─ `pct1`, `pct2` → "95.0%" / "95.00%"
      ─ `kwh`          → "1.234.567 kWh"
      ─ `kwp`          → "1.234,5 kWp"
    """
    if value is None or value == "":
        return "—"
    if fmt is None:
        if isinstance(value, bool):
            return "Sí" if value else "No"
        if isinstance(value, int):
            return fmt_ar(float(value), 0)
        if isinstance(value, float):
            return fmt_ar(value, 0) if value.is_integer() else fmt_ar(value, 2)
        return str(value)

    fmt = fmt.strip().lower()
    if fmt.isdigit():
        return fmt_ar(float(value), int(fmt))
    if fmt == "usd":
        return f"USD {fmt_ar(float(value), 0)}"
    if fmt == "pct":
        return f"{float(value) * 100:.0f}%"
    if fmt.startswith("pct") and fmt[3:].isdigit():
        return f"{float(value) * 100:.{int(fmt[3:])}f}%"
    if fmt == "kwh":
        return f"{fmt_ar(float(value), 0)} kWh"
    if fmt == "kwp":
        return f"{fmt_ar(float(value), 1)} kWp"
    return str(value)


def contexto_de_propuesta(
    *,
    factura,
    sizing,
    inv_cfg,
    costeo,
    cliente_nombre: str = "",
    proyecto_nombre: str = "",
    fecha: datetime.date | None = None,
) -> dict[str, Any]:
    """Arma el contexto que alimenta la sustitución de placeholders.

    Todas las claves de este dict son los identificadores válidos para
    `{{clave}}` en el template.
    """
    f = fecha or datetime.date.today()
    return {
        # Cliente / proyecto
        "cliente": cliente_nombre or factura.titular or "—",
        "titular": factura.titular or "—",
        "proyecto": proyecto_nombre or "—",
        "direccion": factura.direccion or "—",
        "nis": factura.nis or "—",
        "fecha": f.strftime("%d/%m/%Y"),
        "anio": f.year,
        # Consumo eléctrico
        "distribuidora": factura.distribuidora,
        "categoria_tarifaria": factura.categoria_tarifaria,
        "tension": factura.tension_suministro,
        "potencia_contratada": factura.potencia_contratada_kw,
        "consumo_anual": factura.consumo_anual_kwh,
        "consumo_mensual_promedio": factura.consumo_mensual_promedio,
        # Sizing técnico
        "kwp": sizing.kwp_real,
        "n_paneles": sizing.n_paneles,
        "wp_panel": 720,                          # TCL-MG720DT210-68NS (módulo EPC)
        "generacion_anual": sizing.generacion_anual_kwh,
        "cobertura_pct": sizing.cobertura,
        "yield_especifico": (
            sizing.generacion_anual_kwh / sizing.kwp_real if sizing.kwp_real else 0
        ),
        "topologia": getattr(sizing, "topologia", ""),
        # Equipamiento
        "n_inversores": inv_cfg.cantidad,
        "inversor_sku": inv_cfg.inversor.sku,
        "inversor_descripcion": inv_cfg.inversor.descripcion,
        # Inversión (vista cliente — sin márgenes desglosados)
        "neto_usd": costeo.neto_cliente,
        "iva_usd": costeo.iva_total,
        "total_usd": costeo.total_cliente,
        "usd_kwp": (
            costeo.total_cliente / sizing.kwp_real if sizing.kwp_real else 0
        ),
    }


def listar_placeholders(prs) -> set[str]:
    """Devuelve el conjunto de claves `{{...}}` detectadas en el deck.

    Útil para validar que un template no tenga typos antes de pasarlo a la pipeline.
    """
    keys: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            for txt in _iter_textos(shape):
                for m in _PLACEHOLDER_RE.finditer(txt):
                    keys.add(m.group(1))
    return keys


def listar_chart_markers(prs) -> set[str]:
    """Devuelve el conjunto de marcadores `{{chart:KEY}}` detectados en el deck."""
    keys: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            for txt in _iter_textos(shape):
                for m in _CHART_RE.finditer(txt):
                    keys.add(m.group(1))
    return keys


def tiene_placeholders(prs) -> bool:
    """True si el deck contiene al menos un placeholder `{{...}}`."""
    for slide in prs.slides:
        for shape in slide.shapes:
            for txt in _iter_textos(shape):
                if _PLACEHOLDER_RE.search(txt):
                    return True
    return False


def sustituir(prs, contexto: dict[str, Any]) -> int:
    """Reemplaza in-place todos los placeholders. Devuelve cuántos sustituyó.

    Cuando un placeholder cruza varios `runs` de un párrafo (algo que pasa si
    PowerPoint partió el texto al editarlo), fusionamos el texto en el primer
    run y vaciamos los siguientes. Esto puede aplanar formato heterogéneo
    dentro del párrafo — los placeholders pensados para reemplazo deberían
    ir en un único formato.
    """
    total = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            total += _sustituir_shape(shape, contexto)
    return total


def insertar_charts(prs, *, factura=None, sizing=None) -> int:
    """Encuentra marcadores `{{chart:KEY}}` y los reemplaza por charts reales.

    El marcador es un textbox cualquiera con el texto `{{chart:KEY}}`. El
    motor:
      1. Detecta su posición y tamaño en el slide.
      2. Lo elimina.
      3. Inserta un chart real en la misma posición y tamaño.

    Charts soportados:
      - `consumo_mensual`         → bar chart de factura.consumos (necesita ≥3 meses)
      - `cobertura`               → pie chart "Solar vs Red"
      - `generacion_vs_consumo`   → bar chart 2 series (consumo vs gen estimada)

    Si falta data para un chart, deja el marcador (el cliente verá el
    placeholder y sabrá que faltó info).
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    # Recopilamos primero (no podemos modificar el slide mientras iteramos shapes).
    markers: list[tuple] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = "".join(
                r.text or "" for p in shape.text_frame.paragraphs for r in p.runs
            )
            m = _CHART_RE.search(txt)
            if m:
                markers.append((slide, shape, m.group(1)))

    count = 0
    for slide, shape, key in markers:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        creo_chart = False

        if key == "consumo_mensual" and factura and len(factura.consumos) >= 3:
            data = CategoryChartData()
            data.categories = [c.mes for c in factura.consumos]
            data.add_series("kWh consumidos", [c.kwh_total for c in factura.consumos])
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, data,
            )
            creo_chart = True

        elif key == "cobertura" and sizing:
            cob = max(0.0, min(float(sizing.cobertura), 1.0))
            data = CategoryChartData()
            data.categories = ["Solar", "Red eléctrica"]
            data.add_series("Origen energía", [cob, 1 - cob])
            slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, left, top, width, height, data,
            )
            creo_chart = True

        elif key == "generacion_vs_consumo" and factura and sizing and factura.consumos:
            data = CategoryChartData()
            data.categories = [c.mes for c in factura.consumos]
            data.add_series("Consumo", [c.kwh_total for c in factura.consumos])
            # Distribución de generación: usa hsp_mensual si está, sino reparto plano
            gen_mensual = _generacion_por_mes(sizing, factura.consumos)
            data.add_series("Generación estimada", gen_mensual)
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, data,
            )
            creo_chart = True

        if creo_chart:
            # Eliminar el shape marker — el chart ocupa su lugar.
            shape._element.getparent().remove(shape._element)
            count += 1

    return count


def _generacion_por_mes(sizing, consumos: list) -> list[float]:
    """Distribución estimada de generación anual sobre los meses de la factura.

    Si los meses cubren un año completo, usa un perfil estacional típico
    argentino (verano alto, invierno bajo). Para series parciales o no
    contiguas, reparte plano.
    """
    n = len(consumos) or 1
    promedio = sizing.generacion_anual_kwh / 12.0
    # Perfil normalizado (suma = 12) para hemisferio sur AR.
    perfil_ar = [1.18, 1.12, 1.05, 0.92, 0.78, 0.70, 0.74, 0.86, 0.96, 1.06, 1.15, 1.22]
    salida: list[float] = []
    for c in consumos:
        try:
            mm = int(c.mes.split("-")[1])
            salida.append(promedio * perfil_ar[mm - 1])
        except (ValueError, IndexError):
            salida.append(promedio)
    return salida


# ──────────────────────────────────────────────────────────────────────────────
# Helpers internos
# ──────────────────────────────────────────────────────────────────────────────
def _iter_textos(shape):
    """Yields todos los strings de texto dentro de un shape (text frames + celdas de tabla)."""
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            yield "".join(r.text or "" for r in para.runs)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    yield "".join(r.text or "" for r in para.runs)


def _sustituir_shape(shape, contexto: dict[str, Any]) -> int:
    count = 0
    if shape.has_text_frame:
        count += _sustituir_textframe(shape.text_frame, contexto)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                count += _sustituir_textframe(cell.text_frame, contexto)
    return count


def _sustituir_textframe(tf, contexto: dict[str, Any]) -> int:
    count = 0
    for para in tf.paragraphs:
        if not para.runs:
            continue
        full = "".join(r.text or "" for r in para.runs)
        if not _PLACEHOLDER_RE.search(full):
            continue

        def _repl(m: re.Match) -> str:
            nonlocal count
            count += 1
            key, fmt = m.group(1), m.group(2)
            return _formatear(contexto.get(key), fmt)

        nuevo = _PLACEHOLDER_RE.sub(_repl, full)
        para.runs[0].text = nuevo
        for r in para.runs[1:]:
            r.text = ""
    return count
