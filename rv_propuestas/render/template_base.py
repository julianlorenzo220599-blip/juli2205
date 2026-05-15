"""Generador de template `.pptx` base con identidad de marca RV Energía (oct 2025).

Aplica el Manual de Marca RV Energía:

  PALETA PRIMARIA
    Verde lime    #A6FF00  (RGB 166/255/0)   — acento, números, último-palabra titles
    Azul corporativo #1B39CE (RGB 27/57/206) — fondos de sección, énfasis

  PALETA SECUNDARIA
    Negro         #000000   — fondos de portada / cover de sección
    Negro suave   #1A1A1A   — variantes
    Blanco        #FFFFFF   — texto sobre fondos oscuros
    Gris claro    #F2F2F2   — fondos secundarios
    Gris medio    #7A7A7A   — bajadas, body text

  TIPOGRAFÍA
    Outfit (Bold + Light) — títulos combinan Light + última palabra Bold
    Si el cliente no tiene Outfit instalada, PowerPoint hace fallback a Calibri.

  ESTRUCTURA
    Portada → Contenidos → [Sección oscura + Slides claras]× → Cierre
    Logo positivo en slides claras, negativo en slides oscuras/azules.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ──────────────────────────────────────────────────────────────────────────────
# IDENTIDAD DE MARCA RV (Manual oct 2025)
# ──────────────────────────────────────────────────────────────────────────────
VERDE_LIME = RGBColor(0xA6, 0xFF, 0x00)
AZUL = RGBColor(0x1B, 0x39, 0xCE)
NEGRO = RGBColor(0x00, 0x00, 0x00)
NEGRO_SUAVE = RGBColor(0x1A, 0x1A, 0x1A)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_CLARO = RGBColor(0xF2, 0xF2, 0xF2)
GRIS_MEDIO = RGBColor(0x7A, 0x7A, 0x7A)
GRIS_OSCURO = RGBColor(0x33, 0x33, 0x33)

FUENTE = "Outfit"

_ASSETS = Path(__file__).parent / "assets"
LOGO_POSITIVO = _ASSETS / "logo_rv_positivo.png"   # para slides claras
LOGO_NEGATIVO = _ASSETS / "logo_rv_negativo.png"   # para slides oscuras/azules


def crear_template_base(output_path: str | Path) -> Path:
    """Genera el `.pptx` corporativo siguiendo el Manual de Marca RV oct/2025."""
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

    _slide_portada(prs.slides.add_slide(blank))
    _slide_contenidos(prs.slides.add_slide(blank))
    _slide_seccion(prs.slides.add_slide(blank), "01", "Análisis del", "consumo.")
    _slide_analisis_consumo(prs.slides.add_slide(blank))
    _slide_historico_consumo(prs.slides.add_slide(blank))
    _slide_seccion(prs.slides.add_slide(blank), "02", "Solución", "técnica.")
    _slide_solucion(prs.slides.add_slide(blank))
    _slide_generacion_vs_consumo(prs.slides.add_slide(blank))
    _slide_seccion(prs.slides.add_slide(blank), "03", "Inversión y", "próximos pasos.")
    _slide_inversion(prs.slides.add_slide(blank))
    _slide_proximos_pasos(prs.slides.add_slide(blank))

    prs.save(str(output_path))
    return output_path


# ──────────────────────────────────────────────────────────────────────────────
# SLIDES
# ──────────────────────────────────────────────────────────────────────────────
def _slide_portada(slide):
    _fondo(slide, NEGRO)
    _logo(slide, "negativo", left=0.5, top=0.4, height=0.9)
    _chip(slide, "2025", left=0.5, top=2.6)
    # Título estilo manual: "Propuesta técnica." con palabra final en lime
    _titulo_dual(
        slide,
        "Propuesta", "técnica.",
        left=0.5, top=3.2,
        size=80, color1=BLANCO, color2=VERDE_LIME,
    )
    _texto(slide, "{{kwp|kwp}}  ·  {{cliente}}",
           left=0.5, top=5.6, width=12, size=22, color=BLANCO, bold=False)
    _texto(slide, "Proyecto: {{proyecto}}  ·  {{fecha}}",
           left=0.5, top=6.2, width=12, size=14, color=GRIS_MEDIO)

    # Glow azul radial decorativo (rectángulo con tinte)
    _rectangulo(slide, left=9.0, top=0, width=4.5, height=4.5,
                fill=AZUL, line=AZUL, alpha=True)


def _slide_contenidos(slide):
    _fondo(slide, AZUL)
    _logo(slide, "negativo", left=11.5, top=0.4, height=0.7)
    _titulo(slide, "Contenidos.", left=0.5, top=0.6, color=BLANCO, size=52)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.7), Inches(1.0), Inches(0.06),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = VERDE_LIME
    bar.line.fill.background()

    items = [
        ("01", "Análisis del consumo"),
        ("02", "Solución técnica propuesta"),
        ("03", "Generación esperada"),
        ("04", "Inversión"),
        ("05", "Próximos pasos"),
    ]
    for i, (num, txt) in enumerate(items):
        y = 2.6 + i * 0.7
        _texto(slide, num, left=0.7, top=y, width=0.8, size=22,
               bold=True, color=VERDE_LIME)
        _texto(slide, "/", left=1.5, top=y + 0.05, width=0.3, size=20,
               color=BLANCO)
        _texto(slide, txt, left=2.0, top=y, width=10, size=20,
               color=BLANCO)


def _slide_seccion(slide, numero: str, parte1: str, parte2: str):
    """Slide separadora de sección estilo '01 Logotipo.' del manual."""
    _fondo(slide, NEGRO)
    _logo(slide, "negativo", left=0.5, top=0.4, height=0.9)
    _chip(slide, "2025", left=0.5, top=2.6)
    _texto(slide, numero, left=0.5, top=3.4, width=2.5,
           size=96, bold=True, color=VERDE_LIME)
    _titulo_dual(
        slide, parte1, parte2,
        left=3.0, top=3.6,
        size=72, color1=BLANCO, color2=BLANCO, bold2=True,
    )


def _slide_analisis_consumo(slide):
    _fondo(slide, BLANCO)
    _logo(slide, "positivo", left=11.5, top=0.4, height=0.7)
    _titulo_dual(slide, "Análisis del", "consumo.",
                 left=0.5, top=0.5, size=38,
                 color1=GRIS_OSCURO, color2=AZUL, bold2=True)
    _barra_inferior(slide)

    _kpi(slide, 0.5, 1.8, "Cliente",        "{{cliente}}")
    _kpi(slide, 3.8, 1.8, "Distribuidora",  "{{distribuidora}}")
    _kpi(slide, 7.1, 1.8, "Categoría",      "{{categoria_tarifaria}}")
    _kpi(slide, 10.4, 1.8, "NIS",           "{{nis}}", size_value=18)

    _kpi(slide, 0.5, 3.8, "Consumo anual",       "{{consumo_anual|kwh}}",
         accent=VERDE_LIME)
    _kpi(slide, 3.8, 3.8, "Promedio mensual",    "{{consumo_mensual_promedio|kwh}}")
    _kpi(slide, 7.1, 3.8, "Tensión",             "{{tension}}")
    _kpi(slide, 10.4, 3.8, "Pot. contratada",    "{{potencia_contratada|0}} kW")

    _texto(slide, "Domicilio del suministro: {{direccion}}",
           left=0.5, top=6.0, width=12, size=12, color=GRIS_MEDIO)


def _slide_historico_consumo(slide):
    _fondo(slide, BLANCO)
    _logo(slide, "positivo", left=11.5, top=0.4, height=0.7)
    _titulo_dual(slide, "Histórico", "mensual.",
                 left=0.5, top=0.5, size=38,
                 color1=GRIS_OSCURO, color2=AZUL, bold2=True)
    _barra_inferior(slide)
    _texto(slide,
           "Consumos extraídos directamente de la factura del cliente.",
           left=0.5, top=1.5, width=12, size=14, color=GRIS_MEDIO)
    _chart_marker(slide, "consumo_mensual",
                  left=0.8, top=2.1, width=11.7, height=4.6)


def _slide_solucion(slide):
    _fondo(slide, BLANCO)
    _logo(slide, "positivo", left=11.5, top=0.4, height=0.7)
    _titulo_dual(slide, "Solución técnica", "propuesta.",
                 left=0.5, top=0.5, size=36,
                 color1=GRIS_OSCURO, color2=AZUL, bold2=True)
    _barra_inferior(slide)

    _kpi(slide, 0.5, 1.8, "Potencia DC",      "{{kwp|kwp}}", accent=AZUL)
    _kpi(slide, 3.8, 1.8, "Generación anual", "{{generacion_anual|kwh}}",
         accent=VERDE_LIME)
    _kpi(slide, 7.1, 1.8, "Cobertura",        "{{cobertura_pct|pct1}}")
    _kpi(slide, 10.4, 1.8, "Yield",           "{{yield_especifico|0}} kWh/kWp")

    _texto(slide, "Equipamiento principal",
           left=0.5, top=3.9, width=12, size=18, bold=True, color=AZUL)
    bullets = [
        "{{n_paneles}} × módulos TCL TOPCon Bifacial {{wp_panel}} Wp",
        "{{n_inversores}} × {{inversor_descripcion}}",
        "Smart meter trifásico GoodWe + estructura + balance of system",
        "Ingeniería, gestión ante {{distribuidora}}, puesta en marcha",
    ]
    for i, b in enumerate(bullets):
        # Punto verde + texto
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.55), Inches(4.5 + i * 0.45), Inches(0.15), Inches(0.15),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = VERDE_LIME
        dot.line.fill.background()
        _texto(slide, b, left=0.85, top=4.42 + i * 0.45,
               width=12, size=14, color=GRIS_OSCURO)

    _texto(slide,
           "Garantías: 15 años producto · 30 años performance lineal · 10 años inversores",
           left=0.5, top=6.5, width=12, size=11, color=GRIS_MEDIO)


def _slide_generacion_vs_consumo(slide):
    _fondo(slide, BLANCO)
    _logo(slide, "positivo", left=11.5, top=0.4, height=0.7)
    _titulo_dual(slide, "Generación vs", "consumo.",
                 left=0.5, top=0.5, size=36,
                 color1=GRIS_OSCURO, color2=AZUL, bold2=True)
    _barra_inferior(slide)
    _texto(slide,
           "Cobertura de {{cobertura_pct|pct1}} del consumo anual ({{consumo_anual|kwh}}) "
           "con generación estimada de {{generacion_anual|kwh}}.",
           left=0.5, top=1.6, width=12.5, size=13, color=GRIS_MEDIO)
    _chart_marker(slide, "generacion_vs_consumo",
                  left=0.5, top=2.3, width=8.5, height=4.5)
    _chart_marker(slide, "cobertura",
                  left=9.3, top=2.3, width=3.5, height=4.5)


def _slide_inversion(slide):
    _fondo(slide, BLANCO)
    _logo(slide, "positivo", left=11.5, top=0.4, height=0.7)
    _titulo_dual(slide, "Inversión llave", "en mano.",
                 left=0.5, top=0.5, size=36,
                 color1=GRIS_OSCURO, color2=AZUL, bold2=True)
    _barra_inferior(slide)

    _kpi(slide, 0.5, 1.8, "Subtotal (sin IVA)",   "{{neto_usd|usd}}")
    _kpi(slide, 3.8, 1.8, "IVA diferencial",      "{{iva_usd|usd}}")
    _kpi(slide, 7.1, 1.8, "TOTAL",                "{{total_usd|usd}}",
         accent=VERDE_LIME, size_value=24)
    _kpi(slide, 10.4, 1.8, "USD por kWp",         "{{usd_kwp|0}}")

    _texto(slide, "Condiciones comerciales",
           left=0.5, top=3.9, width=12, size=18, bold=True, color=AZUL)
    cond = [
        "Validez de oferta: 30 días corridos",
        "Forma de pago: 50% anticipo · 40% acopio · 10% PEM (negociable)",
        "Plazo de ejecución: 60-90 días desde aprobación de proyecto",
        "Incluye: ingeniería, materiales, instalación, gestión, puesta en marcha",
        "Excluye: obra civil mayor y trámites municipales (cotizable aparte)",
    ]
    for i, c in enumerate(cond):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.55), Inches(4.5 + i * 0.4), Inches(0.13), Inches(0.13),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = AZUL
        dot.line.fill.background()
        _texto(slide, c, left=0.85, top=4.42 + i * 0.4,
               width=12, size=13, color=GRIS_OSCURO)


def _slide_proximos_pasos(slide):
    _fondo(slide, AZUL)
    _logo(slide, "negativo", left=11.5, top=0.4, height=0.7)
    _titulo(slide, "Próximos pasos.", left=0.5, top=0.6,
            color=BLANCO, size=48)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.6), Inches(1.0), Inches(0.06),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = VERDE_LIME
    bar.line.fill.background()

    pasos = [
        ("01", "Visita técnica al sitio y relevamiento del PDI"),
        ("02", "Ingeniería de detalle + simulación PVSyst"),
        ("03", "Gestión de habilitación ante {{distribuidora}}"),
        ("04", "Provisión de equipos y montaje"),
        ("05", "Puesta en marcha, capacitación y entrega"),
    ]
    for i, (num, txt) in enumerate(pasos):
        y = 2.3 + i * 0.65
        _texto(slide, num, left=0.7, top=y, width=0.8, size=22,
               bold=True, color=VERDE_LIME)
        _texto(slide, txt, left=1.7, top=y, width=11, size=17,
               color=BLANCO)

    # Contacto al pie
    _texto(slide, "Contacto",
           left=0.5, top=6.0, width=4, size=14, bold=True, color=VERDE_LIME)
    _texto(slide, "Julián Lorenzo  ·  Project Manager — Nuevos Negocios",
           left=0.5, top=6.4, width=12, size=12, color=BLANCO)
    _texto(slide, "julian.lorenzo@radiovictoria.com.ar  ·  +54 11 4407-6575",
           left=0.5, top=6.7, width=12, size=12, color=BLANCO)


# ──────────────────────────────────────────────────────────────────────────────
# HELPERS VISUALES (alineados con el manual de marca)
# ──────────────────────────────────────────────────────────────────────────────
def _fondo(slide, color: RGBColor) -> None:
    """Pinta el fondo del slide con un rectángulo full-bleed."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    # Mover al fondo
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def _logo(slide, variante: str, *, left: float, top: float, height: float) -> None:
    """Inserta el logo de RV. `variante` = "positivo" (fondos claros) o "negativo"."""
    path = LOGO_POSITIVO if variante == "positivo" else LOGO_NEGATIVO
    if path.exists():
        slide.shapes.add_picture(
            str(path), Inches(left), Inches(top), height=Inches(height),
        )


def _chip(slide, texto: str, *, left: float, top: float) -> None:
    """Píldora con borde verde — usada en portadas (estilo manual)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(0.9), Inches(0.4),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NEGRO
    shape.line.color.rgb = VERDE_LIME
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.margin_top = Inches(0.05)
    tf.text = texto
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _aplicar_fuente(p.runs[0], size=14, color=VERDE_LIME, bold=False)


def _titulo(slide, texto: str, *, left: float, top: float,
            color: RGBColor = AZUL, size: int = 40) -> None:
    _texto(slide, texto, left=left, top=top, width=12.5,
           size=size, bold=True, color=color, height=0.9)


def _titulo_dual(slide, parte1: str, parte2: str, *, left: float, top: float,
                 size: int = 40, color1: RGBColor = GRIS_OSCURO,
                 color2: RGBColor = AZUL, bold2: bool = True) -> None:
    """Título con dos pesos/colores: 'Propuesta técnica.' donde 'técnica.' va resaltado."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(12.5), Inches(size * 0.025 + 0.4),
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    # Run 1 — Light, color1
    r1 = p.add_run()
    r1.text = parte1 + " "
    _aplicar_fuente(r1, size=size, color=color1, bold=False)
    # Run 2 — Bold, color2 (palabra acento)
    r2 = p.add_run()
    r2.text = parte2
    _aplicar_fuente(r2, size=size, color=color2, bold=bold2)


def _texto(slide, texto: str, *, left: float, top: float, width: float = 6.0,
           height: float = 0.5, size: int = 14, bold: bool = False,
           color: RGBColor = GRIS_OSCURO) -> None:
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = texto
    _aplicar_fuente(r, size=size, color=color, bold=bold)


def _kpi(slide, left: float, top: float, label: str, value: str,
         *, size_value: int = 20, accent: RGBColor = AZUL) -> None:
    """KPI tile estilo manual: caja blanca con borde sutil + label gris + valor en accent."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(3.0), Inches(1.7),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRIS_CLARO
    shape.line.color.rgb = accent
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.25)
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = label
    r1 = p1.runs[0]
    _aplicar_fuente(r1, size=10, color=GRIS_MEDIO, bold=True)

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = value
    _aplicar_fuente(r2, size=size_value, color=accent, bold=True)


def _chart_marker(slide, key: str, *, left: float, top: float,
                  width: float, height: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRIS_CLARO
    shape.line.color.rgb = GRIS_MEDIO
    shape.line.dash_style = 7  # dashed

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(height / 2 - 0.25)
    tf.text = "{{chart:" + key + "}}"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _aplicar_fuente(p.runs[0], size=13, color=GRIS_MEDIO, bold=True)


def _barra_inferior(slide) -> None:
    """Barra horizontal azul al pie (decorativa, estilo manual)."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(7.2), Inches(1.2), Inches(0.08),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = AZUL
    bar.line.fill.background()


def _rectangulo(slide, *, left: float, top: float, width: float, height: float,
                fill: RGBColor, line: RGBColor, alpha: bool = False) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0)
    if alpha:
        # Opacidad 40% via XML — simulando efecto glow del manual
        from pptx.oxml.ns import qn
        sp = shape.fill.fore_color._xFill
        srgb = sp.find(qn("a:srgbClr"))
        if srgb is not None:
            alpha_el = srgb.makeelement(qn("a:alpha"), {"val": "40000"})
            srgb.append(alpha_el)


def _aplicar_fuente(run, *, size: int, color: RGBColor, bold: bool) -> None:
    run.font.name = FUENTE
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
