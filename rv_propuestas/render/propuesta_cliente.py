"""Renderizador de propuesta para el cliente (PPT).

Genera un PowerPoint con python-pptx (cross-platform) basado en el template
de RV Energía. Si no hay template, genera uno limpio desde cero.

Vista cliente: NO incluye márgenes ni desglose interno. Solo total final + IVA.
"""
from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from ..config import fmt_ar
from ..costeo.calculo import CosteoResult
from ..inputs.facturas import Factura
from ..inputs.ubicacion import IrradiacionAnual
from ..sizing.engine import SizingResult
from ..sizing.topologia import ConfiguracionInversores


AZUL_RV = RGBColor(0x1F, 0x4E, 0x78)
GRIS = RGBColor(0x59, 0x59, 0x59)
NARANJA = RGBColor(0xED, 0x7D, 0x31)


def _add_title(slide, text: str, top: float = 0.3):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.5), Inches(0.8))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_RV


def _add_bullet(slide, text: str, left: float, top: float, width: float = 6.0, size: int = 14, bold: bool = False, color: RGBColor = GRIS):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.5))
    tf = box.text_frame
    tf.text = text
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color


def _add_kpi_box(slide, left: float, top: float, label: str, value: str, color: RGBColor = AZUL_RV):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(3.0), Inches(1.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
    shape.line.color.rgb = color
    tf = shape.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)

    tf.text = label
    p1 = tf.paragraphs[0]
    p1.font.size = Pt(12)
    p1.font.color.rgb = GRIS

    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = color


def render(
    factura: Factura,
    irradiacion: IrradiacionAnual,
    sizing: SizingResult,
    inv_cfg: ConfiguracionInversores,
    costeo: CosteoResult,
    output_path: str | Path,
    cliente_nombre: str = "",
    proyecto_nombre: str = "",
    template_path: Optional[str | Path] = None,
) -> Path:
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    if template_path and Path(template_path).exists():
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

    # ── Slide 1: Portada ──────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    _add_title(s1, "Propuesta Técnico-Económica", top=2.0)
    _add_bullet(s1, f"Planta solar fotovoltaica {fmt_ar(sizing.kwp_real, 0)} kWp",
                left=0.5, top=3.0, width=12, size=22, bold=True, color=NARANJA)
    _add_bullet(s1, f"Cliente: {cliente_nombre or factura.titular or '—'}",
                left=0.5, top=4.0, width=12, size=18)
    _add_bullet(s1, f"Proyecto: {proyecto_nombre or '—'}",
                left=0.5, top=4.5, width=12, size=18)
    _add_bullet(s1, "RV Energía  ·  Radio Victoria S.A.",
                left=0.5, top=6.7, width=12, size=12, color=AZUL_RV, bold=True)

    # ── Slide 2: Datos del cliente ────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    _add_title(s2, "Análisis del Consumo")
    _add_bullet(s2, f"Distribuidora: {factura.distribuidora}", left=0.5, top=1.5, size=16)
    _add_bullet(s2, f"Categoría tarifaria: {factura.categoria_tarifaria}", left=0.5, top=2.0, size=16)
    _add_bullet(s2, f"Tensión: {factura.tension_suministro}", left=0.5, top=2.5, size=16)
    _add_bullet(s2, f"Potencia contratada: {factura.potencia_contratada_kw or '—'} kW", left=0.5, top=3.0, size=16)

    _add_kpi_box(s2, 7.0, 1.5, "Consumo anual", f"{fmt_ar(factura.consumo_anual_kwh, 0)} kWh")
    _add_kpi_box(s2, 7.0, 3.4, "Consumo mensual prom.", f"{fmt_ar(factura.consumo_mensual_promedio, 0)} kWh")

    # ── Slide 3: Solución propuesta ───────────────────────────────────────────
    s3 = prs.slides.add_slide(blank)
    _add_title(s3, "Solución Propuesta")
    _add_kpi_box(s3, 0.5, 1.5, "Potencia DC", f"{fmt_ar(sizing.kwp_real, 1)} kWp", color=NARANJA)
    _add_kpi_box(s3, 3.7, 1.5, "Generación anual", f"{fmt_ar(sizing.generacion_anual_kwh, 0)} kWh")
    _add_kpi_box(s3, 6.9, 1.5, "Cobertura consumo", f"{sizing.cobertura * 100:.0f}%")
    _add_kpi_box(s3, 10.1, 1.5, "Paneles", f"{sizing.n_paneles} × 725 W")

    _add_bullet(s3, "Equipamiento principal:", left=0.5, top=3.6, size=16, bold=True, color=AZUL_RV)
    _add_bullet(s3, f"• Módulos TCL TOPCon Bifacial 725 Wp — {sizing.n_paneles} unidades",
                left=0.7, top=4.1, width=12, size=14)
    _add_bullet(s3, f"• Inversores GoodWe: {inv_cfg.cantidad} × {inv_cfg.inversor.descripcion}",
                left=0.7, top=4.5, width=12, size=14)
    _add_bullet(s3, f"• Estructura metálica + obra eléctrica completa",
                left=0.7, top=4.9, width=12, size=14)
    _add_bullet(s3, f"• Ingeniería, gestión ante distribuidora, puesta en marcha",
                left=0.7, top=5.3, width=12, size=14)

    _add_bullet(s3, f"Garantías: Módulos 15 años producto / 30 años performance lineal · Inversores 10 años",
                left=0.5, top=6.3, width=12, size=12, color=GRIS)

    # ── Slide 4: Inversión ────────────────────────────────────────────────────
    s4 = prs.slides.add_slide(blank)
    _add_title(s4, "Inversión")

    _add_bullet(s4, "Inversión total llave en mano:", left=0.5, top=1.5, size=18, bold=True, color=AZUL_RV)

    _add_kpi_box(s4, 0.5, 2.3, "Subtotal (sin IVA)", f"USD {fmt_ar(costeo.neto_cliente, 0)}")
    _add_kpi_box(s4, 3.7, 2.3, "IVA", f"USD {fmt_ar(costeo.iva_total, 0)}")
    _add_kpi_box(s4, 6.9, 2.3, "TOTAL (USD)", f"USD {fmt_ar(costeo.total_cliente, 0)}", color=NARANJA)
    _add_kpi_box(s4, 10.1, 2.3, "USD por kWp", f"{fmt_ar(costeo.total_cliente / sizing.kwp_real, 0)}")

    _add_bullet(s4, "Condiciones comerciales:", left=0.5, top=4.5, size=16, bold=True, color=AZUL_RV)
    _add_bullet(s4, "• Validez de oferta: 30 días", left=0.7, top=5.0, width=12, size=14)
    _add_bullet(s4, "• Forma de pago: a convenir (anticipo 50% / 40% acopio / 10% PEM)", left=0.7, top=5.4, width=12, size=14)
    _add_bullet(s4, "• Plazo de ejecución: 60-90 días desde aprobación de proyecto", left=0.7, top=5.8, width=12, size=14)
    _add_bullet(s4, "• Incluye: ingeniería, materiales, instalación, gestión, puesta en marcha", left=0.7, top=6.2, width=12, size=14)

    # ── Slide 5: Próximos pasos ───────────────────────────────────────────────
    s5 = prs.slides.add_slide(blank)
    _add_title(s5, "Próximos Pasos")
    _add_bullet(s5, "1. Visita técnica al sitio y relevamiento del PDI", left=0.7, top=1.8, width=12, size=16)
    _add_bullet(s5, "2. Ingeniería de detalle + simulación PVSyst", left=0.7, top=2.5, width=12, size=16)
    _add_bullet(s5, "3. Gestión de habilitación ante distribuidora", left=0.7, top=3.2, width=12, size=16)
    _add_bullet(s5, "4. Provisión de equipos y montaje", left=0.7, top=3.9, width=12, size=16)
    _add_bullet(s5, "5. Puesta en marcha + capacitación + entrega final", left=0.7, top=4.6, width=12, size=16)

    _add_bullet(s5, "Contacto:", left=0.5, top=5.8, size=14, bold=True, color=AZUL_RV)
    _add_bullet(s5, "Julián Lorenzo — Project Manager, Nuevos Negocios", left=0.5, top=6.2, width=12, size=12)
    _add_bullet(s5, "julian.lorenzo@radiovictoria.com.ar  ·  +54 11 4407-6575", left=0.5, top=6.5, width=12, size=12)

    prs.save(str(output_path))
    return output_path
