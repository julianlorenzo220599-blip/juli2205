"""Tests del template engine (sustitución de placeholders en .pptx).

Construye templates sintéticos en memoria con python-pptx, los procesa con
el motor de sustitución y verifica el texto final. No requiere archivos
.pptx externos.
"""
from __future__ import annotations

from io import BytesIO
from pathlib import Path
from types import SimpleNamespace

from pptx import Presentation
from pptx.util import Inches

from rv_propuestas.render.template import (
    contexto_de_propuesta,
    insertar_charts,
    listar_chart_markers,
    listar_placeholders,
    sustituir,
    tiene_placeholders,
)
from rv_propuestas.render.template import _formatear  # privado pero útil para tests
from rv_propuestas.render.template_base import crear_template_base


# ──────────────────────────────────────────────────────────────────────────────
# Helpers para construir templates de prueba
# ──────────────────────────────────────────────────────────────────────────────
def _template_con_textos(*lineas: str) -> Presentation:
    """Crea un .pptx en memoria con N slides, cada una con un textbox simple."""
    prs = Presentation()
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    for txt in lineas:
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tb.text_frame.text = txt
    return prs


def _texto_de_slide(prs, idx: int) -> str:
    s = prs.slides[idx]
    out = []
    for shape in s.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                out.append("".join(r.text or "" for r in para.runs))
    return "\n".join(out)


# ──────────────────────────────────────────────────────────────────────────────
# Filtros de formato
# ──────────────────────────────────────────────────────────────────────────────
def test_formatear_sin_filtro_entero():
    assert _formatear(1234, None) == "1.234"
    assert _formatear(1234.0, None) == "1.234"
    assert _formatear(12.5, None) == "12,50"


def test_formatear_decimales():
    assert _formatear(1234.567, "0") == "1.235"
    assert _formatear(1234.567, "1") == "1.234,6"
    assert _formatear(1234.567, "2") == "1.234,57"


def test_formatear_usd():
    assert _formatear(1234567.89, "usd") == "USD 1.234.568"


def test_formatear_pct():
    assert _formatear(0.95, "pct") == "95%"
    assert _formatear(0.953, "pct1") == "95.3%"


def test_formatear_kwh_kwp():
    assert _formatear(556332, "kwh") == "556.332 kWh"
    assert _formatear(242.4, "kwp") == "242,4 kWp"


def test_formatear_none():
    assert _formatear(None, None) == "—"
    assert _formatear(None, "usd") == "—"


# ──────────────────────────────────────────────────────────────────────────────
# Detección
# ──────────────────────────────────────────────────────────────────────────────
def test_tiene_placeholders_true():
    prs = _template_con_textos("Hola {{cliente}}")
    assert tiene_placeholders(prs)


def test_tiene_placeholders_false():
    prs = _template_con_textos("Sin placeholders aquí")
    assert not tiene_placeholders(prs)


def test_listar_placeholders():
    prs = _template_con_textos(
        "Cliente: {{cliente}}",
        "Potencia: {{kwp|1}} kWp · Inversión {{total_usd|usd}}",
    )
    ks = listar_placeholders(prs)
    assert ks == {"cliente", "kwp", "total_usd"}


# ──────────────────────────────────────────────────────────────────────────────
# Sustitución
# ──────────────────────────────────────────────────────────────────────────────
def test_sustituir_basico():
    prs = _template_con_textos("Cliente: {{cliente}} — {{proyecto}}")
    ctx = {"cliente": "ACME SA", "proyecto": "Planta 250 kW"}
    n = sustituir(prs, ctx)
    assert n == 2
    assert _texto_de_slide(prs, 0) == "Cliente: ACME SA — Planta 250 kW"


def test_sustituir_con_filtros():
    prs = _template_con_textos(
        "Potencia {{kwp|1}} kWp · Inversión {{total_usd|usd}} · Cobertura {{cobertura|pct}}"
    )
    ctx = {"kwp": 242.42, "total_usd": 145073.5, "cobertura": 0.951}
    sustituir(prs, ctx)
    assert _texto_de_slide(prs, 0) == "Potencia 242,4 kWp · Inversión USD 145.074 · Cobertura 95%"


def test_sustituir_clave_inexistente():
    prs = _template_con_textos("Falta: {{no_existe}}")
    sustituir(prs, {"otra": "x"})
    assert _texto_de_slide(prs, 0) == "Falta: —"


def test_sustituir_persiste_en_pptx():
    """Roundtrip: sustituir, guardar a bytes, releer, verificar."""
    prs = _template_con_textos("Total: {{total_usd|usd}}")
    sustituir(prs, {"total_usd": 1234567.89})
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    reloaded = Presentation(buf)
    assert _texto_de_slide(reloaded, 0) == "Total: USD 1.234.568"


# ──────────────────────────────────────────────────────────────────────────────
# Contexto desde una propuesta
# ──────────────────────────────────────────────────────────────────────────────
def test_contexto_de_propuesta_completo():
    factura = SimpleNamespace(
        titular="ACME SA", distribuidora="EDESUR", categoria_tarifaria="T3-MT",
        tension_suministro="MT 13.2 kV", potencia_contratada_kw=200,
        consumo_anual_kwh=480000, consumo_mensual_promedio=40000,
        direccion="Av. Test 123", nis="80515768",
    )
    sizing = SimpleNamespace(
        kwp_real=242.4, n_paneles=334,
        generacion_anual_kwh=399548, cobertura=0.951,
    )
    inv = SimpleNamespace(
        cantidad=1, inversor=SimpleNamespace(sku="GW225K-HTH", descripcion="GoodWe 225kW"),
    )
    costeo = SimpleNamespace(
        neto_cliente=120000, iva_total=25073, total_cliente=145073,
    )
    ctx = contexto_de_propuesta(
        factura=factura, sizing=sizing, inv_cfg=inv, costeo=costeo,
        cliente_nombre="ACME SA", proyecto_nombre="Planta 250 kW",
    )

    assert ctx["cliente"] == "ACME SA"
    assert ctx["proyecto"] == "Planta 250 kW"
    assert ctx["distribuidora"] == "EDESUR"
    assert ctx["kwp"] == 242.4
    assert ctx["total_usd"] == 145073
    assert abs(ctx["usd_kwp"] - 145073 / 242.4) < 1e-6
    assert ctx["nis"] == "80515768"


def test_contexto_renderiza_en_template():
    """Test integrado: contexto + sustitución + verificación de texto resultante."""
    factura = SimpleNamespace(
        titular="ACME", distribuidora="EDESUR", categoria_tarifaria="T3-MT",
        tension_suministro="MT", potencia_contratada_kw=200,
        consumo_anual_kwh=480000, consumo_mensual_promedio=40000,
        direccion="—", nis="123",
    )
    sizing = SimpleNamespace(
        kwp_real=242.4, n_paneles=334, generacion_anual_kwh=399548, cobertura=0.951,
    )
    inv = SimpleNamespace(cantidad=1, inversor=SimpleNamespace(sku="GW225K-HTH", descripcion="GoodWe 225kW"))
    costeo = SimpleNamespace(neto_cliente=120000, iva_total=25073, total_cliente=145073)
    ctx = contexto_de_propuesta(
        factura=factura, sizing=sizing, inv_cfg=inv, costeo=costeo,
        cliente_nombre="ACME", proyecto_nombre="Planta 250 kW",
    )

    prs = _template_con_textos(
        "Cliente: {{cliente}}",
        "Solución: {{kwp|kwp}} con {{n_paneles}} paneles {{wp_panel}} Wp",
        "Inversión total: {{total_usd|usd}} (USD/kWp: {{usd_kwp|0}})",
    )
    sustituir(prs, ctx)

    assert _texto_de_slide(prs, 0) == "Cliente: ACME"
    assert _texto_de_slide(prs, 1) == "Solución: 242,4 kWp con 334 paneles 720 Wp"
    assert "USD 145.073" in _texto_de_slide(prs, 2)


# ──────────────────────────────────────────────────────────────────────────────
# Template base + chart injection
# ──────────────────────────────────────────────────────────────────────────────
def _template_con_chart_marker(key: str) -> Presentation:
    """Crea un .pptx con un textbox que contiene `{{chart:KEY}}` como marcador."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    tb.text_frame.text = "{{chart:" + key + "}}"
    return prs


def test_listar_chart_markers_detecta_tipos():
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for key in ("consumo_mensual", "cobertura", "generacion_vs_consumo"):
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(3))
        tb.text_frame.text = "{{chart:" + key + "}}"
    detectados = listar_chart_markers(prs)
    assert detectados == {"consumo_mensual", "cobertura", "generacion_vs_consumo"}


def test_listar_placeholders_ignora_chart_markers():
    """Los marcadores `{{chart:X}}` no deben aparecer en `listar_placeholders` —
    son una sintaxis aparte."""
    prs = _template_con_chart_marker("consumo_mensual")
    assert listar_placeholders(prs) == set()


def test_insertar_chart_consumo_mensual():
    prs = _template_con_chart_marker("consumo_mensual")
    factura = SimpleNamespace(
        consumos=[
            SimpleNamespace(mes=f"2025-{m:02d}", kwh_total=1000 + m * 100)
            for m in range(1, 13)
        ],
    )
    n = insertar_charts(prs, factura=factura, sizing=None)
    assert n == 1
    # Verificar que se agregó un chart al slide
    slide = prs.slides[0]
    charts = [sh for sh in slide.shapes if sh.has_chart]
    assert len(charts) == 1
    # Verificar que el marker desapareció
    textos = [
        "".join(r.text or "" for r in p.runs)
        for sh in slide.shapes if sh.has_text_frame
        for p in sh.text_frame.paragraphs
    ]
    assert not any("chart:consumo_mensual" in t for t in textos)


def test_insertar_chart_cobertura_pie():
    prs = _template_con_chart_marker("cobertura")
    sizing = SimpleNamespace(cobertura=0.85, generacion_anual_kwh=100, kwp_real=10)
    n = insertar_charts(prs, factura=None, sizing=sizing)
    assert n == 1
    slide = prs.slides[0]
    assert sum(1 for sh in slide.shapes if sh.has_chart) == 1


def test_insertar_chart_consumo_mensual_skip_si_pocos_meses():
    """Con <3 meses no tiene sentido un histórico — debe dejar el marker."""
    prs = _template_con_chart_marker("consumo_mensual")
    factura = SimpleNamespace(
        consumos=[SimpleNamespace(mes="2025-01", kwh_total=1000)],
    )
    n = insertar_charts(prs, factura=factura, sizing=None)
    assert n == 0
    slide = prs.slides[0]
    assert sum(1 for sh in slide.shapes if sh.has_chart) == 0


def test_insertar_chart_sin_datos_no_crashea():
    """Si no se pasa factura/sizing, no debe inyectar nada y no debe lanzar."""
    prs = _template_con_chart_marker("consumo_mensual")
    n = insertar_charts(prs, factura=None, sizing=None)
    assert n == 0


def test_insertar_chart_generacion_vs_consumo():
    prs = _template_con_chart_marker("generacion_vs_consumo")
    factura = SimpleNamespace(
        consumos=[
            SimpleNamespace(mes=f"2025-{m:02d}", kwh_total=40000)
            for m in range(1, 13)
        ],
    )
    sizing = SimpleNamespace(
        generacion_anual_kwh=456000, kwp_real=300, cobertura=0.95,
    )
    n = insertar_charts(prs, factura=factura, sizing=sizing)
    assert n == 1


def test_crear_template_base(tmp_path=None):
    import tempfile
    out = Path(tempfile.mkdtemp()) / "tb.pptx"
    crear_template_base(out)
    assert out.exists() and out.stat().st_size > 1000
    prs = Presentation(out)
    # 8 slides minimalistas: portada, contenidos, análisis, histórico,
    # solución, gen vs consumo, inversión, próximos pasos
    assert len(prs.slides) == 8
    # Placeholders esperados
    ph = listar_placeholders(prs)
    assert {"cliente", "kwp", "total_usd", "cobertura_pct"}.issubset(ph)
    # 3 chart markers
    assert listar_chart_markers(prs) == {
        "consumo_mensual", "cobertura", "generacion_vs_consumo",
    }


def test_contexto_incluye_yield_y_topologia():
    """contexto_de_propuesta debe exponer yield_especifico y topologia."""
    factura = SimpleNamespace(
        titular="X", distribuidora="EDESUR", categoria_tarifaria="T3",
        tension_suministro="MT", potencia_contratada_kw=100,
        consumo_anual_kwh=100000, consumo_mensual_promedio=8333,
        direccion="—", nis="X",
    )
    sizing = SimpleNamespace(
        kwp_real=100, n_paneles=140, generacion_anual_kwh=150000,
        cobertura=0.95, topologia="MT_TRAFO",
    )
    inv = SimpleNamespace(cantidad=1, inversor=SimpleNamespace(sku="X", descripcion="X"))
    costeo = SimpleNamespace(neto_cliente=100, iva_total=20, total_cliente=120)
    ctx = contexto_de_propuesta(
        factura=factura, sizing=sizing, inv_cfg=inv, costeo=costeo,
    )
    assert ctx["yield_especifico"] == 1500.0
    assert ctx["topologia"] == "MT_TRAFO"


if __name__ == "__main__":
    import sys
    fns = [v for k, v in sorted(globals().items()) if k.startswith("test_") and callable(v)]
    fail = 0
    for fn in fns:
        try:
            fn()
            print(f"✓ {fn.__name__}")
        except AssertionError as e:
            print(f"✗ {fn.__name__}: {e}")
            fail += 1
        except Exception as e:
            print(f"✗ {fn.__name__}: {type(e).__name__}: {e}")
            fail += 1
    sys.exit(0 if fail == 0 else 1)
