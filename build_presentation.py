"""Build improved RV Energía residential proposal presentation."""
from pathlib import Path
from copy import deepcopy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsmap

ROOT = Path(__file__).parent
ASSETS = ROOT / "assets"
OUT = ROOT / "output" / "RV_Energia_Propuesta_Residencial_Country_2026_v2.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# Slide size — 16:9 widescreen 13.333" x 7.5"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT = "Outfit"

# Palette
LIGHT_BG = RGBColor(0xFA, 0xFA, 0xF7)   # warm off-white
LIGHT_INK = RGBColor(0x12, 0x14, 0x18)
LIGHT_INK_SOFT = RGBColor(0x55, 0x5A, 0x60)
LIGHT_LINE = RGBColor(0xE3, 0xE3, 0xDE)
DARK_BG = RGBColor(0x0B, 0x0E, 0x14)
DARK_INK = RGBColor(0xF5, 0xF5, 0xF0)
DARK_INK_SOFT = RGBColor(0xA0, 0xA6, 0xAE)
DARK_LINE = RGBColor(0x23, 0x2A, 0x34)
ACCENT_GREEN = RGBColor(0xAE, 0xFF, 0x00)   # logo lime
ACCENT_BLUE = RGBColor(0x29, 0x52, 0xD3)    # logo blue
GOLD = RGBColor(0xE8, 0xB5, 0x4A)


def set_font(run, name=FONT, size=None, bold=None, color=None):
    run.font.name = name
    # also set the East Asian + Symbol typeface so PowerPoint keeps Outfit
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
        el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", name)
    latin = rPr.find(qn("a:latin"))
    if latin is None:
        latin = etree.SubElement(rPr, qn("a:latin"))
    latin.set("typeface", name)
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line=1.15, tracking=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        r = p.add_run()
        r.text = ln
        set_font(r, size=size, bold=bold, color=color)
        if tracking:
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(tracking))
    return tb


def add_rich(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line=1.15):
    """runs: list of (text, dict_of_props). New paragraph triggered by text containing only '\\n'."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line
    for text, props in runs:
        if text == "\n":
            p = tf.add_paragraph()
            p.alignment = props.get("align", align)
            p.line_spacing = props.get("line", line)
            continue
        r = p.add_run()
        r.text = text
        set_font(r,
                 size=props.get("size", 14),
                 bold=props.get("bold", False),
                 color=props.get("color"))
        if props.get("tracking"):
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(props["tracking"]))
    return tb


def fill_solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, rgb, *, line_rgb=None, line_w=None,
             corner=None):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp_type, x, y, w, h)
    s.shadow.inherit = False
    if corner is not None:
        try:
            s.adjustments[0] = corner
        except Exception:
            pass
    if rgb is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = rgb
    if line_rgb is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_rgb
        if line_w is not None:
            s.line.width = line_w
    return s


def add_hline(slide, x, y, w, rgb, weight=Emu(6350)):
    ln = slide.shapes.add_connector(1, x, y, x + w, y)
    ln.line.color.rgb = rgb
    ln.line.width = weight
    return ln


def set_slide_bg(slide, rgb):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_picture_with_transparency(slide, path, x, y, w, h, alpha_pct=100):
    """Insert picture and apply alpha (0-100, where 100 = fully opaque)."""
    pic = slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    if alpha_pct < 100:
        # inject <a:alphaModFix amt="..."/> inside the blip
        blip = pic._element.find(".//" + qn("a:blip"))
        # remove existing alpha modifiers
        for tag in ("a:alphaModFix", "a:alphaMod"):
            ex = blip.find(qn(tag))
            if ex is not None:
                blip.remove(ex)
        amf = etree.SubElement(blip, qn("a:alphaModFix"))
        amf.set("amt", str(int(alpha_pct * 1000)))
    return pic


def add_logo(slide, dark=False, x=None, y=None, w=Inches(1.1)):
    path = ASSETS / ("logo_white.png" if dark else "logo_color.png")
    # logo aspect: white = 858x454 (~1.89), color = 1296x405 (~3.2)
    aspect = 858/454 if dark else 1296/405
    if x is None:
        x = Inches(0.5)
    if y is None:
        y = Inches(0.45)
    h = Emu(int(w / aspect))
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def page_footer(slide, n, dark=False, chapter=""):
    ink = DARK_INK_SOFT if dark else LIGHT_INK_SOFT
    line = DARK_LINE if dark else LIGHT_LINE
    y = Inches(7.05)
    add_hline(slide, Inches(0.6), y, Inches(12.13), line, weight=Emu(3175))
    add_text(slide, Inches(0.6), y + Inches(0.05), Inches(8), Inches(0.3),
             chapter.upper(), size=8, color=ink, tracking=200)
    add_text(slide, Inches(11.7), y + Inches(0.05), Inches(1.1), Inches(0.3),
             f"{n:02d} / 12", size=8, color=ink, align=PP_ALIGN.RIGHT, tracking=200)


def add_eyebrow(slide, x, y, text, dark=False, color=None):
    c = color or (ACCENT_GREEN if dark else ACCENT_BLUE)
    add_text(slide, x, y, Inches(6), Inches(0.3), text, size=9, bold=True,
             color=c, tracking=400)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ===========================================================
    # SLIDE 1 — COVER (DARK · hero photo at ~20% opacity)
    # ===========================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, DARK_BG)
    # Hero photo as full-bleed background at 20% opacity (80% transparent)
    add_picture_with_transparency(s, ASSETS / "hero_house.jpg",
                                  0, 0, SLIDE_W, SLIDE_H, alpha_pct=20)
    # Soft dark vignette to ensure legibility on bottom-left
    add_rect(s, 0, Inches(3.0), Inches(8.5), Inches(4.5), DARK_BG)
    # actually make a gradient using two overlapping rectangles with alpha
    # (kept simple: solid background already dark)

    add_logo(s, dark=True, x=Inches(0.7), y=Inches(0.55), w=Inches(1.4))

    # Top-right meta
    add_text(s, Inches(9.3), Inches(0.65), Inches(3.5), Inches(0.3),
             "PROPUESTA RESIDENCIAL · 2026", size=9, bold=True,
             color=ACCENT_GREEN, align=PP_ALIGN.RIGHT, tracking=400)
    add_text(s, Inches(9.3), Inches(0.95), Inches(3.5), Inches(0.3),
             "EDICIÓN COUNTRY · ZONA NORTE GBA", size=8,
             color=DARK_INK_SOFT, align=PP_ALIGN.RIGHT, tracking=300)

    # Big headline
    add_rich(s, Inches(0.7), Inches(3.0), Inches(11), Inches(2.6), [
        ("Tu casa.", {"size": 80, "bold": True, "color": DARK_INK, "line": 1.0}),
        ("\n", {}),
        ("Tu propia central ", {"size": 80, "bold": True, "color": DARK_INK, "line": 1.0}),
        ("de energía.", {"size": 80, "bold": True, "color": ACCENT_GREEN, "line": 1.0}),
    ])
    add_text(s, Inches(0.7), Inches(5.95), Inches(9.5), Inches(0.8),
             "Propuesta integral llave en mano para incorporar energía solar fotovoltaica\n"
             "en tu vivienda — pensada para quien busca tranquilidad, eficiencia y un activo que se paga solo.",
             size=12, color=DARK_INK_SOFT, line=1.5)

    # Bottom chips
    chips = ["RV ENERGÍA", "78 AÑOS DE HISTORIA", "TIER 1 BLOOMBERGNEF", "EDICIÓN COUNTRY"]
    cx = Inches(0.7)
    cy = Inches(6.95)
    for i, c in enumerate(chips):
        w = Inches(0.04 + 0.085*len(c))
        r = add_rect(s, cx, cy, w, Inches(0.32), None,
                     line_rgb=DARK_LINE, line_w=Emu(6350), corner=0.5)
        tb = add_text(s, cx, cy, w, Inches(0.32), c, size=8, bold=True,
                      color=DARK_INK, align=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE, tracking=300)
        cx += w + Inches(0.15)

    # ===========================================================
    # SLIDE 2 — QUIÉNES SOMOS (LIGHT)
    # ===========================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT_BG)
    add_logo(s, dark=False, x=Inches(0.6), y=Inches(0.45), w=Inches(1.4))
    add_eyebrow(s, Inches(0.6), Inches(1.5), "CAPÍTULO 01 · QUIÉNES SOMOS")
    add_rich(s, Inches(0.6), Inches(1.85), Inches(8), Inches(2.0), [
        ("Radio Victoria,", {"size": 56, "bold": True, "color": LIGHT_INK, "line": 1.0}),
        ("\n", {}),
        ("ahora ", {"size": 56, "bold": True, "color": LIGHT_INK, "line": 1.0}),
        ("solar.", {"size": 56, "bold": True, "color": ACCENT_BLUE, "line": 1.0}),
    ])
    add_text(s, Inches(0.6), Inches(4.0), Inches(7.5), Inches(0.4),
             "78 AÑOS HACIENDO INDUSTRIA", size=10, bold=True,
             color=LIGHT_INK_SOFT, tracking=400)
    add_text(s, Inches(0.6), Inches(4.4), Inches(7.5), Inches(2.0),
             "RV Energía es la división solar de Radio Victoria Argentina S.A. — "
             "una compañía nacional con casi ocho décadas en electrónica, logística "
             "y operaciones técnicas en todo el país. Hoy operamos en tres frentes: "
             "distribución mayorista a instaladores, proyectos EPC llave en mano "
             "para industria y agro, y kits residenciales para clientes finales — "
             "respaldados por logística federal e ingeniería propia.",
             size=12, color=LIGHT_INK_SOFT, line=1.55)

    # Right-side photo (solar panels with sky)
    add_picture_with_transparency(s, ASSETS / "panel_sky.jpg",
                                  Inches(8.7), Inches(1.5),
                                  Inches(4.1), Inches(5.0))

    # Stats row
    stats = [("1947", "FUNDACIÓN · RADIO VICTORIA"),
             ("33.000 m²", "OPERACIÓN INDUSTRIAL FEDERAL"),
             ("24/7", "MONITOREO REMOTO · SOPORTE")]
    sx = Inches(0.6); sy = Inches(6.5)
    for big, lbl in stats:
        add_text(s, sx, sy, Inches(2.5), Inches(0.5), big,
                 size=28, bold=True, color=ACCENT_BLUE)
        add_text(s, sx, sy + Inches(0.55), Inches(2.5), Inches(0.3),
                 lbl, size=8, color=LIGHT_INK_SOFT, tracking=300)
        sx += Inches(2.7)

    page_footer(s, 2, dark=False, chapter="RV Energía · Radio Victoria Argentina")

    # ===========================================================
    # SLIDE 3 — ALIADOS TECNOLÓGICOS (LIGHT)
    # ===========================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT_BG)
    add_logo(s, dark=False, x=Inches(0.6), y=Inches(0.45), w=Inches(1.4))
    add_eyebrow(s, Inches(0.6), Inches(1.5), "CAPÍTULO 02 · ALIADOS TECNOLÓGICOS")
    add_rich(s, Inches(0.6), Inches(1.85), Inches(11.5), Inches(1.6), [
        ("Top mundial en cada ", {"size": 48, "bold": True, "color": LIGHT_INK, "line": 1.0}),
        ("componente.", {"size": 48, "bold": True, "color": ACCENT_BLUE, "line": 1.0}),
    ])
    add_text(s, Inches(0.6), Inches(3.2), Inches(11.5), Inches(0.7),
             "Trabajamos exclusivamente con fabricantes Tier 1 BloombergNEF — la clasificación financiera "
             "más alta del mercado solar global, reservada a los proveedores con mayor solvencia y trayectoria.",
             size=12, color=LIGHT_INK_SOFT, line=1.45)

    # Two big cards
    card_specs = [
        {"title": "TCL Solar",
         "sub": "PANELES FOTOVOLTAICOS · CHINA",
         "desc": "Uno de los mayores fabricantes del mundo. Tecnología N-Type TOPCon Bifacial — más eficiencia por m², más vida útil.",
         "rows": [("Panel", "585 Wp · TOPCon Bifacial"),
                  ("Garantía", "15 años producto · 30 potencia"),
                  ("Resistencia", "IP68 · −40 °C a +85 °C")]},
        {"title": "GoodWe",
         "sub": "INVERSORES Y BATERÍAS · CHINA",
         "desc": "Líder global en inversores solares e híbridos. Monitoreo SEMS+ 24/7 vía app + baterías LiFePO4.",
         "rows": [("Inversores", "XS · MS · ES · ET · SDT"),
                  ("Baterías", "Lynx U (LV) · Lynx D (HV)"),
                  ("Garantía", "10 años (extensible)")]}
    ]
    cx = Inches(0.6); cy = Inches(4.15); cw = Inches(6.05); ch = Inches(2.7)
    for spec in card_specs:
        add_rect(s, cx, cy, cw, ch, None, line_rgb=LIGHT_LINE,
                 line_w=Emu(6350), corner=0.04)
        add_text(s, cx + Inches(0.35), cy + Inches(0.3), Inches(5), Inches(0.5),
                 spec["title"], size=22, bold=True, color=LIGHT_INK)
        add_text(s, cx + Inches(0.35), cy + Inches(0.75), Inches(5), Inches(0.3),
                 spec["sub"], size=8, color=ACCENT_BLUE, bold=True, tracking=300)
        add_text(s, cx + Inches(0.35), cy + Inches(1.1), Inches(5.3), Inches(0.7),
                 spec["desc"], size=10, color=LIGHT_INK_SOFT, line=1.4)
        ry = cy + Inches(1.9)
        for k, v in spec["rows"]:
            add_text(s, cx + Inches(0.35), ry, Inches(1.7), Inches(0.25),
                     k.upper(), size=7, color=LIGHT_INK_SOFT, tracking=300)
            add_text(s, cx + Inches(2.0), ry, Inches(3.8), Inches(0.25),
                     v, size=10, bold=True, color=LIGHT_INK)
            ry += Inches(0.27)
        cx += cw + Inches(0.25)

    page_footer(s, 3, dark=False, chapter="Capítulo 02 · Aliados Tecnológicos")

    # ===========================================================
    # SLIDE 4 — CÓMO FUNCIONA (DARK · contrast)
    # ===========================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, DARK_BG)
    # subtle photo on right at low opacity
    add_picture_with_transparency(s, ASSETS / "panel_field.jpg",
                                  Inches(7.5), 0, Inches(5.83), SLIDE_H,
                                  alpha_pct=30)
    add_logo(s, dark=True, x=Inches(0.6), y=Inches(0.45), w=Inches(1.4))
    add_eyebrow(s, Inches(0.6), Inches(1.5), "CAPÍTULO 03 · CÓMO FUNCIONA",
                dark=True)
    add_rich(s, Inches(0.6), Inches(1.85), Inches(11.5), Inches(1.6), [
        ("El recorrido de la ", {"size": 48, "bold": True, "color": DARK_INK, "line": 1.0}),
        ("energía.", {"size": 48, "bold": True, "color": ACCENT_GREEN, "line": 1.0}),
    ])
    add_text(s, Inches(0.6), Inches(3.2), Inches(11.5), Inches(0.4),
             "Del techo al enchufe, en menos de un segundo.",
             size=13, color=DARK_INK_SOFT)

    steps = [("01", "Sol", "FUENTE", "Una radiación gratis, todos los días, todo el año."),
             ("02", "Panel", "DC", "Convierte la luz en corriente continua."),
             ("03", "Inversor", "DC → AC", "Pasa DC a AC 220 V — la misma de tu enchufe."),
             ("04", "Tu casa", "CONSUMO + RED", "Lo sobrante se inyecta a la red y descuenta tu factura.")]
    sx = Inches(0.6); sy = Inches(3.95); sw = Inches(2.95); sh = Inches(2.1)
    for i, (n, t, k, d) in enumerate(steps):
        x = sx + i*(sw + Inches(0.12))
        add_rect(s, x, sy, sw, sh, None, line_rgb=DARK_LINE,
                 line_w=Emu(6350), corner=0.05)
        add_text(s, x + Inches(0.3), sy + Inches(0.25), Inches(1), Inches(0.3),
                 n, size=10, bold=True, color=ACCENT_GREEN, tracking=200)
        add_text(s, x + Inches(0.3), sy + Inches(0.6), Inches(2.5), Inches(0.5),
                 t, size=22, bold=True, color=DARK_INK)
        add_text(s, x + Inches(0.3), sy + Inches(1.1), Inches(2.5), Inches(0.25),
                 k, size=8, bold=True, color=DARK_INK_SOFT, tracking=300)
        add_text(s, x + Inches(0.3), sy + Inches(1.4), Inches(2.5), Inches(0.7),
                 d, size=9, color=DARK_INK_SOFT, line=1.35)
        # arrow between
        if i < 3:
            ax = x + sw + Inches(0.005)
            add_text(s, ax, sy + Inches(0.85), Inches(0.12), Inches(0.4),
                     "→", size=18, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

    # Bottom 3 scenarios
    scenarios = [("MIENTRAS HAY SOL", "Generás y consumís.",
                  "Los electrodomésticos toman energía directo del techo. Cero costo en boleta."),
                 ("SI GENERÁS DE MÁS", "Inyectás a la red.",
                  "El excedente vuelve a la red. La distribuidora te lo descuenta."),
                 ("DE NOCHE O NUBLADO", "Tomás de red o batería.",
                  "Consumís de la red, o de tu batería si el kit es híbrido.")]
    sy2 = Inches(6.2); sw2 = Inches(3.99)
    for i, (k, t, d) in enumerate(scenarios):
        x = Inches(0.6) + i*(sw2 + Inches(0.1))
        add_text(s, x, sy2, sw2, Inches(0.25),
                 k, size=8, bold=True, color=ACCENT_GREEN, tracking=300)
        add_text(s, x, sy2 + Inches(0.25), sw2, Inches(0.35),
                 t, size=14, bold=True, color=DARK_INK)
        add_text(s, x, sy2 + Inches(0.6), sw2, Inches(0.4),
                 d, size=9, color=DARK_INK_SOFT, line=1.35)

    page_footer(s, 4, dark=True, chapter="Capítulo 03 · Cómo Funciona")

    # ===========================================================
    # SLIDE 5 — COMPONENTES (LIGHT)
    # ===========================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT_BG)
    add_logo(s, dark=False, x=Inches(0.6), y=Inches(0.45), w=Inches(1.4))
    add_eyebrow(s, Inches(0.6), Inches(1.5), "CAPÍTULO 04 · COMPONENTES")
    add_rich(s, Inches(0.6), Inches(1.85), Inches(11.5), Inches(1.6), [
        ("¿Qué hay detrás ", {"size": 46, "bold": True, "color": LIGHT_INK, "line": 1.0}),
        ("del sistema?", {"size": 46, "bold": True, "color": ACCENT_BLUE, "line": 1.0}),
    ])
    add_text(s, Inches(0.6), Inches(3.15), Inches(11), Inches(0.4),
             "Cinco partes. Una misma ingeniería.",
             size=12, color=LIGHT_INK_SOFT)

    parts = [("01", "Panel solar", "TCL · 585 Wp TOPCon",
              "Capta la luz del sol y la convierte en electricidad. Va sobre el techo, en estructura de acero."),
             ("02", "Inversor", "GoodWe · Mono/Trifásico",
              "El cerebro del sistema. Convierte DC en AC y administra el flujo entre techo, red y batería."),
             ("03", "Batería (opcional)", "LiFePO4 · Lynx U / Lynx D",
              "Guarda energía para usar de noche o si se corta la luz. Solo en kits híbridos."),
             ("04", "Smart Meter", "Smart Energy Controller",
              "Mide cuánto generás, consumís e inyectás. La caja registradora del sistema."),
             ("05", "Estructura + cables", "Steel Frame · 1500 V DC",
              "Soportes de acero, cables, conectores y protecciones eléctricas.")]
    cx = Inches(0.6); cy = Inches(3.85); cw = Inches(2.41); ch = Inches(2.4)
    for i, (n, t, k, d) in enumerate(parts):
        x = cx + i*(cw + Inches(0.1))
        add_rect(s, x, cy, cw, ch, None, line_rgb=LIGHT_LINE,
                 line_w=Emu(6350), corner=0.05)
        add_text(s, x + Inches(0.25), cy + Inches(0.25), cw, Inches(0.3),
                 n, size=9, bold=True, color=ACCENT_BLUE, tracking=300)
        add_text(s, x + Inches(0.25), cy + Inches(0.55), cw - Inches(0.3), Inches(0.7),
                 t, size=15, bold=True, color=LIGHT_INK, line=1.1)
        add_text(s, x + Inches(0.25), cy + Inches(1.25), cw - Inches(0.3), Inches(0.3),
                 k.upper(), size=7, bold=True, color=LIGHT_INK_SOFT, tracking=200)
        add_text(s, x + Inches(0.25), cy + Inches(1.55), cw - Inches(0.3), Inches(0.85),
                 d, size=9, color=LIGHT_INK_SOFT, line=1.35)

    # Bonus banner with SEMS+
    bx = Inches(0.6); by = Inches(6.45); bw = Inches(12.13); bh = Inches(0.55)
    add_rect(s, bx, by, bw, bh, LIGHT_INK, corner=0.3)
    add_text(s, bx + Inches(0.4), by, Inches(11), bh,
             "BONUS · App SEMS+ en tu celular · 24/7 · sin costo · incluida en todos los kits",
             size=11, bold=True, color=LIGHT_BG, anchor=MSO_ANCHOR.MIDDLE, tracking=200)

    page_footer(s, 5, dark=False, chapter="Capítulo 04 · Componentes")

    # ===========================================================
    # SLIDE 6 — ARQUITECTURAS (LIGHT)
    # ===========================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT_BG)
    add_logo(s, dark=False, x=Inches(0.6), y=Inches(0.45), w=Inches(1.4))
    add_eyebrow(s, Inches(0.6), Inches(1.5), "CAPÍTULO 05 · ARQUITECTURAS")
    add_rich(s, Inches(0.6), Inches(1.85), Inches(11.5), Inches(1.6), [
        ("Tres formas de conectar con el ", {"size": 44, "bold": True, "color": LIGHT_INK, "line": 1.0}),
        ("sol.", {"size": 44, "bold": True, "color": ACCENT_BLUE, "line": 1.0}),
    ])
    add_text(s, Inches(0.6), Inches(3.1), Inches(11.5), Inches(0.4),
             "Esta es la decisión más importante del proyecto.",
             size=12, color=LIGHT_INK_SOFT)

    archs = [
        {"badge": "ON-GRID · SIN BATERÍA", "title": "Autoconsumo puro.",
         "desc": "Conectado a la red, le vendés los excedentes. Sin baterías — la opción más rentable, payback más corto.",
         "rows": [("IDEAL SI", "Red estable"),
                  ("CUBRE", "Consumo diurno"),
                  ("RESPALDO", "No")],
         "highlight": False},
        {"badge": "HÍBRIDO · RED + BATERÍA", "title": "Autoconsumo con respaldo.",
         "desc": "Red + batería LiFePO4. Si se corta la luz, tu casa sigue funcionando. La opción más completa.",
         "rows": [("IDEAL SI", "Cortes frecuentes"),
                  ("CUBRE", "Día + noche + cortes"),
                  ("RESPALDO", "Sí · batería 5 kWh")],
         "highlight": True},
        {"badge": "OFF-GRID · SIN RED", "title": "Independencia total.",
         "desc": "Para zonas rurales sin red. Vivís 100% con sol y batería — sin distribuidora, sin factura.",
         "rows": [("IDEAL SI", "Campo / sin red"),
                  ("CUBRE", "Todo el consumo"),
                  ("RESPALDO", "Sí · batería ampliable")],
         "highlight": False},
    ]
    cx = Inches(0.6); cy = Inches(3.85); cw = Inches(4.04); ch = Inches(3.0)
    for i, a in enumerate(archs):
        x = cx + i*(cw + Inches(0.05))
        if a["highlight"]:
            add_rect(s, x, cy, cw, ch, LIGHT_INK, corner=0.04)
            ink = LIGHT_BG; ink_soft = RGBColor(0xB0, 0xB6, 0xBE); badge_c = ACCENT_GREEN
        else:
            add_rect(s, x, cy, cw, ch, None, line_rgb=LIGHT_LINE,
                     line_w=Emu(6350), corner=0.04)
            ink = LIGHT_INK; ink_soft = LIGHT_INK_SOFT; badge_c = ACCENT_BLUE
        add_text(s, x + Inches(0.3), cy + Inches(0.25), cw, Inches(0.25),
                 a["badge"], size=8, bold=True, color=badge_c, tracking=400)
        add_text(s, x + Inches(0.3), cy + Inches(0.55), cw - Inches(0.4), Inches(0.6),
                 a["title"], size=22, bold=True, color=ink, line=1.05)
        add_text(s, x + Inches(0.3), cy + Inches(1.25), cw - Inches(0.4), Inches(0.9),
                 a["desc"], size=10, color=ink_soft, line=1.4)
        ry = cy + Inches(2.15)
        for k, v in a["rows"]:
            add_text(s, x + Inches(0.3), ry, Inches(1.2), Inches(0.22),
                     k, size=7, bold=True, color=ink_soft, tracking=300)
            add_text(s, x + Inches(1.55), ry, Inches(2.3), Inches(0.22),
                     v, size=10, bold=True, color=ink)
            ry += Inches(0.23)
        if a["highlight"]:
            add_text(s, x + Inches(0.3), cy + ch - Inches(0.05),
                     cw, Inches(0.2),
                     "★ RECOMENDADO COUNTRY", size=7, bold=True,
                     color=ACCENT_GREEN, tracking=400)

    page_footer(s, 6, dark=False, chapter="Capítulo 05 · Arquitecturas")

    # ===========================================================
    # SLIDE 7 — LEY 27.424 (DARK · contrast)
    # ===========================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, DARK_BG)
    add_picture_with_transparency(s, ASSETS / "panels_aerial.jpg",
                                  0, 0, SLIDE_W, SLIDE_H, alpha_pct=15)
    add_logo(s, dark=True, x=Inches(0.6), y=Inches(0.45), w=Inches(1.4))
    add_eyebrow(s, Inches(0.6), Inches(1.5), "CAPÍTULO 06 · LEY 27.424",
                dark=True)
    add_rich(s, Inches(0.6), Inches(1.85), Inches(11.5), Inches(1.7), [
        ("Argentina, desde 2017.\n", {"size": 38, "bold": True, "color": DARK_INK_SOFT, "line": 1.0}),
        ("Vos sos un ", {"size": 56, "bold": True, "color": DARK_INK, "line": 1.0}),
        ("usuario-generador.", {"size": 56, "bold": True, "color": ACCENT_GREEN, "line": 1.0}),
    ])
    add_text(s, Inches(0.6), Inches(3.95), Inches(11.5), Inches(1.0),
             "La ley te permite generar tu propia energía y venderle los excedentes a la red eléctrica. "
             "Se llama balance neto o net-metering: la distribuidora cambia el medidor por uno bidireccional "
             "que cuenta los dos flujos — lo que tomás y lo que inyectás.",
             size=12, color=DARK_INK_SOFT, line=1.5)

    steps3 = [
        ("01", "Medidor bidireccional",
         "Cuenta lo que consumís y lo que generás de más. A fin de mes, factura la diferencia."),
        ("02", "Pagás solo la diferencia",
         "Si generaste más de lo que consumiste, tenés crédito para el mes siguiente."),
        ("03", "RV gestiona el trámite",
         "Nosotros nos encargamos del trámite con tu distribuidora (Edenor, Edesur, EDEN, EDEA…).")
    ]
    sx = Inches(0.6); sy = Inches(5.15); sw = Inches(4.04); sh = Inches(1.5)
    for i, (n, t, d) in enumerate(steps3):
        x = sx + i*(sw + Inches(0.05))
        add_rect(s, x, sy, sw, sh, None, line_rgb=DARK_LINE,
                 line_w=Emu(6350), corner=0.04)
        add_text(s, x + Inches(0.25), sy + Inches(0.18), Inches(1), Inches(0.3),
                 n, size=10, bold=True, color=ACCENT_GREEN, tracking=200)
        add_text(s, x + Inches(0.25), sy + Inches(0.45), sw - Inches(0.4), Inches(0.4),
                 t, size=14, bold=True, color=DARK_INK)
        add_text(s, x + Inches(0.25), sy + Inches(0.85), sw - Inches(0.4), Inches(0.6),
                 d, size=9, color=DARK_INK_SOFT, line=1.4)

    # Quote
    add_text(s, Inches(0.6), Inches(6.75), Inches(12.13), Inches(0.32),
             "“Una instalación on-grid bien dimensionada se paga sola entre 4 y 6 años — "
             "el sistema sigue generando energía gratis 25 años más.”",
             size=11, color=DARK_INK_SOFT)

    page_footer(s, 7, dark=True, chapter="Capítulo 06 · Marco Regulatorio")

    # ===========================================================
    # SLIDE 8 — TU KIT IDEAL (LIGHT)
    # ===========================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT_BG)
    add_logo(s, dark=False, x=Inches(0.6), y=Inches(0.45), w=Inches(1.4))
    add_eyebrow(s, Inches(0.6), Inches(1.5), "CAPÍTULO 07 · TU KIT IDEAL")
    add_rich(s, Inches(0.6), Inches(1.85), Inches(7.5), Inches(1.8), [
        ("Nuestra recomendación\n", {"size": 36, "bold": True, "color": LIGHT_INK, "line": 1.0}),
        ("para tu casa en el ", {"size": 36, "bold": True, "color": LIGHT_INK, "line": 1.0}),
        ("country.", {"size": 36, "bold": True, "color": ACCENT_BLUE, "line": 1.0}),
    ])

    # Right hero image
    add_picture_with_transparency(s, ASSETS / "hero_house.jpg",
                                  Inches(8.3), Inches(1.5),
                                  Inches(4.5), Inches(2.8))

    # Recommendation card
    rx = Inches(0.6); ry = Inches(3.95); rw = Inches(7.5); rh = Inches(3.1)
    add_rect(s, rx, ry, rw, rh, LIGHT_INK, corner=0.03)
    add_text(s, rx + Inches(0.4), ry + Inches(0.3), Inches(6), Inches(0.25),
             "★ RECOMENDADO · ALWAYS CONNECTED PLUS · 5 kW",
             size=8, bold=True, color=ACCENT_GREEN, tracking=400)
    add_text(s, rx + Inches(0.4), ry + Inches(0.6), Inches(6), Inches(0.5),
             "12 paneles · 5 kW · batería 5 kWh",
             size=20, bold=True, color=LIGHT_BG)
    add_text(s, rx + Inches(0.4), ry + Inches(1.15), Inches(6.7), Inches(1.4),
             "Híbrido monofásico con respaldo ante cortes. La configuración estrella para zona norte GBA. "
             "Cubre consumo diurno completo (aire, pileta, electrodomésticos), guarda excedente para la noche "
             "y mantiene tu casa funcionando si se corta la luz.",
             size=10, color=RGBColor(0xB8, 0xBE, 0xC6), line=1.45)
    # Three pills
    pills = [("USD 4.399", "+ IVA · kit completo"),
             ("5–7 años", "PAYBACK ESTIMADO"),
             ("~7.500 kWh/año", "GENERACIÓN")]
    px = rx + Inches(0.4); py = ry + Inches(2.45)
    for big, lbl in pills:
        add_text(s, px, py, Inches(2.4), Inches(0.35),
                 big, size=14, bold=True, color=ACCENT_GREEN)
        add_text(s, px, py + Inches(0.35), Inches(2.4), Inches(0.2),
                 lbl, size=7, color=RGBColor(0xB8, 0xBE, 0xC6), tracking=300)
        px += Inches(2.4)

    # Specs table right
    tx = Inches(8.3); ty = Inches(4.5); tw = Inches(4.5); th = Inches(2.55)
    add_rect(s, tx, ty, tw, th, None, line_rgb=LIGHT_LINE,
             line_w=Emu(6350), corner=0.04)
    add_text(s, tx + Inches(0.3), ty + Inches(0.2), Inches(4), Inches(0.3),
             "ESPECIFICACIONES TÉCNICAS", size=8, bold=True,
             color=ACCENT_BLUE, tracking=400)
    specs = [("POTENCIA AC", "5.000 W"),
             ("PANELES", "12 × TCL 585 Wp"),
             ("INVERSOR", "GoodWe GW5000-ES-20"),
             ("BATERÍA", "Lynx U · 5 kWh LiFePO4"),
             ("ESPACIO EN TECHO", "≈ 31 m²")]
    sy = ty + Inches(0.55)
    for k, v in specs:
        add_text(s, tx + Inches(0.3), sy, Inches(1.8), Inches(0.25),
                 k, size=7, color=LIGHT_INK_SOFT, tracking=300)
        add_text(s, tx + Inches(2.2), sy, Inches(2.2), Inches(0.25),
                 v, size=10, bold=True, color=LIGHT_INK)
        sy += Inches(0.38)

    page_footer(s, 8, dark=False, chapter="Capítulo 07 · Tu Kit Ideal")

    # ===========================================================
    # SLIDE 9 — CATÁLOGO (LIGHT)
    # ===========================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT_BG)
    add_logo(s, dark=False, x=Inches(0.6), y=Inches(0.45), w=Inches(1.4))
    add_eyebrow(s, Inches(0.6), Inches(1.5), "CAPÍTULO 08 · CATÁLOGO")
    add_rich(s, Inches(0.6), Inches(1.85), Inches(11.5), Inches(1.0), [
        ("Catálogo residencial ", {"size": 40, "bold": True, "color": LIGHT_INK, "line": 1.0}),
        ("2026.", {"size": 40, "bold": True, "color": ACCENT_BLUE, "line": 1.0}),
    ])
    add_text(s, Inches(0.6), Inches(2.75), Inches(11.5), Inches(0.5),
             "Precios de referencia del kit completo. Conexión monofásica — el formato más común en country.",
             size=11, color=LIGHT_INK_SOFT)

    # Table
    cols = [("KIT", 2.4), ("TIPO", 1.1), ("POT.", 0.7), ("PANELES", 1.3),
            ("BATERÍA", 1.5), ("IDEAL PARA", 3.2), ("USD + IVA", 1.4)]
    rows = [
        ("Economy Basic", "On-Grid", "3 kW", "6 × 585 W", "—", "Familia chica · <400 kWh/mes", "1.099", False),
        ("Economy Plus", "On-Grid", "5 kW", "10 × 585 W", "—", "Hogar promedio · 400–800 kWh/mes", "1.999", False),
        ("Economy Pro", "On-Grid", "6 kW", "12 × 585 W", "—", "Alto consumo · 800–1.200 kWh/mes", "2.299", False),
        ("Always Connected", "Híbrido", "3 kW", "6 × 585 W", "Lynx U 5 kWh", "Country chico · con cortes", "3.199", False),
        ("Always Connected Plus ★", "Híbrido", "5 kW", "12 × 585 W", "Lynx U 5 kWh", "Recomendado country · alto consumo + respaldo", "4.399", True),
        ("Rural", "Off-Grid", "3 kW", "6 × 585 W", "Lynx U 5 kWh", "Cabaña / fin de semana sin red", "3.199", False),
    ]
    tx = Inches(0.6); ty = Inches(3.5); row_h = Inches(0.45)
    # Header row
    cx = tx
    add_hline(s, tx, ty + row_h, Inches(12.13), LIGHT_INK)
    for label, wi in cols:
        add_text(s, cx + Inches(0.1), ty + Inches(0.12), Inches(wi), Inches(0.3),
                 label, size=7, bold=True, color=LIGHT_INK, tracking=400)
        cx += Inches(wi)

    ry = ty + row_h
    for r in rows:
        kit, tipo, pot, pan, bat, ideal, price, hi = r
        if hi:
            add_rect(s, tx, ry, Inches(12.13), row_h, LIGHT_INK, corner=0.02)
            txt_c = LIGHT_BG; sub_c = ACCENT_GREEN
        else:
            txt_c = LIGHT_INK; sub_c = LIGHT_INK_SOFT
            add_hline(s, tx, ry + row_h, Inches(12.13), LIGHT_LINE,
                      weight=Emu(3175))
        vals = [kit, tipo, pot, pan, bat, ideal, price]
        cx = tx
        for (label, wi), v in zip(cols, vals):
            bold = (label == "KIT") or (label == "USD + IVA")
            add_text(s, cx + Inches(0.1), ry + Inches(0.1), Inches(wi), Inches(0.3),
                     v, size=10, bold=bold, color=txt_c if bold else (sub_c if not hi else RGBColor(0xCD, 0xD3, 0xDB)),
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += Inches(wi)
        ry += row_h

    add_text(s, Inches(0.6), Inches(6.75), Inches(12.13), Inches(0.3),
             "PRECIOS USD + IVA · KIT COMPLETO (PANELES, INVERSOR, BATERÍA SI APLICA, SMART METER Y CABLEADO). "
             "ESTRUCTURA Y MANO DE OBRA COTIZADAS POR SEPARADO SEGÚN RELEVAMIENTO.",
             size=7, color=LIGHT_INK_SOFT, tracking=200)

    page_footer(s, 9, dark=False, chapter="Capítulo 08 · Catálogo")

    # ===========================================================
    # SLIDE 10 — PROCESO (LIGHT)
    # ===========================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT_BG)
    add_logo(s, dark=False, x=Inches(0.6), y=Inches(0.45), w=Inches(1.4))
    add_eyebrow(s, Inches(0.6), Inches(1.5), "CAPÍTULO 09 · PROCESO")
    add_rich(s, Inches(0.6), Inches(1.85), Inches(8.5), Inches(1.6), [
        ("Una firma. ", {"size": 38, "bold": True, "color": LIGHT_INK, "line": 1.0}),
        ("Un Project Manager.", {"size": 38, "bold": True, "color": ACCENT_BLUE, "line": 1.0}),
    ])
    add_text(s, Inches(0.6), Inches(3.1), Inches(8.5), Inches(0.8),
             "Vos tratás con una sola persona. Nosotros gestionamos a TCL, GoodWe, "
             "instaladores y distribuidora — todo bajo nuestro mando.",
             size=12, color=LIGHT_INK_SOFT, line=1.5)

    # Right photo: installer hands
    add_picture_with_transparency(s, ASSETS / "installer_hands.jpg",
                                  Inches(9.4), Inches(1.5),
                                  Inches(3.4), Inches(2.3))

    # 6-step grid (2 rows × 3 cols)
    proc = [
        ("01", "Relevamiento", "Visita técnica al techo. Medición de orientación, sombras y espacios."),
        ("02", "Dimensionamiento", "Análisis de tus últimas 6 boletas y propuesta del kit óptimo."),
        ("03", "Ingeniería", "Memorias de cálculo, planos y trámite con la distribuidora."),
        ("04", "Instalación", "Montaje completo en el techo. Entre 1 y 3 días de obra."),
        ("05", "Puesta en marcha", "Conexión a red, configuración de la app SEMS+ y capacitación."),
        ("06", "Monitoreo 24/7", "Seguimiento remoto y post-venta de por vida desde RV."),
    ]
    gx = Inches(0.6); gy = Inches(3.95); cw = Inches(2.95); ch = Inches(1.25)
    for i, (n, t, d) in enumerate(proc):
        col = i % 3; row = i // 3
        x = gx + col*(cw + Inches(0.13))
        y = gy + row*(ch + Inches(0.15))
        add_rect(s, x, y, cw, ch, None, line_rgb=LIGHT_LINE,
                 line_w=Emu(6350), corner=0.05)
        add_text(s, x + Inches(0.25), y + Inches(0.18), Inches(0.8), Inches(0.25),
                 n, size=9, bold=True, color=ACCENT_BLUE, tracking=200)
        add_text(s, x + Inches(0.85), y + Inches(0.18), cw - Inches(1), Inches(0.3),
                 t, size=12, bold=True, color=LIGHT_INK)
        add_text(s, x + Inches(0.25), y + Inches(0.55), cw - Inches(0.4), Inches(0.65),
                 d, size=9, color=LIGHT_INK_SOFT, line=1.4)

    # Bottom highlights (time + warranties)
    hx = Inches(9.4); hy = Inches(3.95); hw = Inches(3.4)
    add_rect(s, hx, hy, hw, Inches(1.25), ACCENT_BLUE, corner=0.05)
    add_text(s, hx + Inches(0.3), hy + Inches(0.15), hw, Inches(0.3),
             "TIEMPO TOTAL", size=8, bold=True, color=ACCENT_GREEN, tracking=400)
    add_text(s, hx + Inches(0.3), hy + Inches(0.4), hw, Inches(0.55),
             "30 a 60 días", size=22, bold=True, color=LIGHT_BG)
    add_text(s, hx + Inches(0.3), hy + Inches(0.9), hw - Inches(0.5), Inches(0.3),
             "Firma a generación · obra 1–3 días",
             size=9, color=RGBColor(0xCD, 0xD3, 0xDB))

    add_rect(s, hx, hy + Inches(1.4), hw, Inches(1.25), LIGHT_INK, corner=0.05)
    add_text(s, hx + Inches(0.3), hy + Inches(1.55), hw, Inches(0.3),
             "GARANTÍAS", size=8, bold=True, color=ACCENT_GREEN, tracking=400)
    add_text(s, hx + Inches(0.3), hy + Inches(1.8), hw, Inches(0.55),
             "15 / 30 / 10 años", size=22, bold=True, color=LIGHT_BG)
    add_text(s, hx + Inches(0.3), hy + Inches(2.3), hw - Inches(0.5), Inches(0.3),
             "Panel · potencia · inversor · 6.000 ciclos batería",
             size=9, color=RGBColor(0xCD, 0xD3, 0xDB))

    page_footer(s, 10, dark=False, chapter="Capítulo 09 · Proceso")

    # ===========================================================
    # SLIDE 11 — FAQ (LIGHT)
    # ===========================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT_BG)
    add_logo(s, dark=False, x=Inches(0.6), y=Inches(0.45), w=Inches(1.4))
    add_eyebrow(s, Inches(0.6), Inches(1.5), "CAPÍTULO 10 · PREGUNTAS FRECUENTES")
    add_rich(s, Inches(0.6), Inches(1.85), Inches(11.5), Inches(1.6), [
        ("Las dudas que aparecen ", {"size": 40, "bold": True, "color": LIGHT_INK, "line": 1.0}),
        ("siempre.", {"size": 40, "bold": True, "color": ACCENT_BLUE, "line": 1.0}),
    ])

    faqs = [
        ("¿Y si está nublado o llueve?",
         "El panel sigue generando — solo que menos. En días nublados produce entre 10% y 30% de su potencia nominal."),
        ("¿De noche cómo hago?",
         "Tomás de la red (on-grid) o de la batería (híbrido). Lo que generaste de día y no usaste te lo descuentan."),
        ("¿Cuánto dura el sistema?",
         "30 años de garantía de potencia en panel, 15 de producto, 10 en inversor (extensible), ~6.000 ciclos en batería."),
        ("¿En cuánto se paga?",
         "En country zona norte, con Always Connected Plus y consumo medio-alto: entre 5 y 7 años."),
        ("¿Es seguro tener eso en el techo?",
         "Sí. Paneles IP68, operan de −40 °C a +85 °C, soportan granizo de hasta 25 mm. Estructura certificada."),
        ("¿Hay que romper el techo?",
         "No. La estructura se fija con anclajes específicos según el techo (teja, chapa, losa). Sin filtraciones."),
        ("¿Necesito permiso del country?",
         "RV prepara la memoria técnica + renders para presentar al consorcio, y el permiso de la distribuidora."),
        ("¿Y el mantenimiento?",
         "Prácticamente cero: limpieza anual y revisión visual. SEMS+ monitorea solo y avisa si algo falla.")
    ]
    gx = Inches(0.6); gy = Inches(3.5); cw = Inches(2.95); ch = Inches(1.65)
    for i, (q, a) in enumerate(faqs):
        col = i % 4; row = i // 4
        x = gx + col*(cw + Inches(0.13))
        y = gy + row*(ch + Inches(0.15))
        add_rect(s, x, y, cw, ch, None, line_rgb=LIGHT_LINE,
                 line_w=Emu(6350), corner=0.05)
        add_text(s, x + Inches(0.25), y + Inches(0.18), Inches(0.5), Inches(0.3),
                 f"{i+1:02d}", size=9, bold=True, color=ACCENT_BLUE, tracking=200)
        add_text(s, x + Inches(0.25), y + Inches(0.45), cw - Inches(0.4), Inches(0.5),
                 q, size=11, bold=True, color=LIGHT_INK, line=1.2)
        add_text(s, x + Inches(0.25), y + Inches(0.95), cw - Inches(0.4), Inches(0.65),
                 a, size=8, color=LIGHT_INK_SOFT, line=1.4)

    page_footer(s, 11, dark=False, chapter="Capítulo 10 · Preguntas Frecuentes")

    # ===========================================================
    # SLIDE 12 — CIERRE (DARK · cinematic)
    # ===========================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, DARK_BG)
    add_picture_with_transparency(s, ASSETS / "panel_sky.jpg",
                                  0, 0, SLIDE_W, SLIDE_H, alpha_pct=25)
    add_logo(s, dark=True, x=Inches(0.6), y=Inches(0.45), w=Inches(1.4))

    add_text(s, Inches(0.6), Inches(1.5), Inches(8), Inches(0.3),
             "CIERRE · HABLEMOS", size=9, bold=True,
             color=ACCENT_GREEN, tracking=400)

    add_rich(s, Inches(0.6), Inches(1.95), Inches(11.5), Inches(2.3), [
        ("El sol está ", {"size": 64, "bold": True, "color": DARK_INK, "line": 1.0}),
        ("gratis.", {"size": 64, "bold": True, "color": ACCENT_GREEN, "line": 1.0}),
        ("\n", {}),
        ("La energía, no. ", {"size": 64, "bold": True, "color": DARK_INK, "line": 1.0}),
        ("Ahorrá.", {"size": 64, "bold": True, "color": ACCENT_GREEN, "line": 1.0}),
    ])

    add_text(s, Inches(0.6), Inches(4.5), Inches(11), Inches(0.8),
             "Coordinamos una visita técnica a tu casa para medir techo, orientación y consumo. "
             "Volvemos con una propuesta detallada, sin compromiso, en menos de 5 días hábiles.",
             size=12, color=DARK_INK_SOFT, line=1.5)

    # Contact blocks
    contact = [
        ("TELÉFONO", ["+54 9 11 4407-6575", "+54 9 11 5140-6628"]),
        ("EMAIL", ["julian.lorenzo@radiovictoria.com.ar",
                   "santiago.basombrio@radiovictoria.com.ar"]),
        ("OFICINAS · WEB", ["Maipú 1210 · Piso 4 · CABA",
                            "www.rvenergia.com.ar"]),
    ]
    cx = Inches(0.6); cy = Inches(5.7)
    for title, lines in contact:
        add_text(s, cx, cy, Inches(4), Inches(0.25),
                 title, size=8, bold=True, color=ACCENT_GREEN, tracking=400)
        for j, ln in enumerate(lines):
            add_text(s, cx, cy + Inches(0.3) + Inches(0.28*j),
                     Inches(4), Inches(0.3),
                     ln, size=11, bold=True, color=DARK_INK)
        cx += Inches(4.2)

    add_text(s, Inches(0.6), Inches(7.1), Inches(8), Inches(0.3),
             "ENERGÍA LIMPIA · UN FUTURO BRILLANTE",
             size=8, bold=True, color=DARK_INK_SOFT, tracking=400)
    add_text(s, Inches(11.7), Inches(7.1), Inches(1.1), Inches(0.3),
             "12 / 12", size=8, color=DARK_INK_SOFT,
             align=PP_ALIGN.RIGHT, tracking=200)

    prs.save(str(OUT))
    print(f"Saved → {OUT}")


if __name__ == "__main__":
    build()
